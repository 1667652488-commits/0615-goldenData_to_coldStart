#!/usr/bin/env python3
"""生成 冷启动规则挖掘工具 PPT v3
配色方案对齐 评估器业界洞察及方案架构.pptx（浅色风格）
章节结构：业界洞察 → 工具介绍 → 案例说明 → 展望
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ==================== 颜色体系 (对齐参考PPT浅色风格) ====================
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)        # 主背景白
BG_LIGHT = RGBColor(0xF8, 0xF9, 0xF3)        # 浅米色背景
BG_LIGHT_GRAY = RGBColor(0xEC, 0xEC, 0xE5)  # 浅灰背景

# 强调色
ACCENT_GREEN = RGBColor(0xE5, 0xF2, 0xEC)    # 浅绿卡片
ACCENT_PURPLE = RGBColor(0xEB, 0xEA, 0xFD)   # 浅紫卡片
ACCENT_PINK = RGBColor(0xFC, 0xF0, 0xF4)     # 浅粉卡片
ACCENT_LAVENDER = RGBColor(0xF1, 0xF1, 0xFD) # 薰衣草卡片
ACCENT_ORANGE_BG = RGBColor(0xF3, 0xC4, 0xB1) # 橙色背景
ACCENT_GOLD = RGBColor(0xFF, 0xC0, 0x00)     # 金黄
ACCENT_LIME = RGBColor(0x92, 0xD0, 0x50)     # 亮绿

# 文字色
TEXT_DARK = RGBColor(0x19, 0x1B, 0x1F)        # 主文字深色
TEXT_GRAY = RGBColor(0x67, 0x67, 0x68)        # 灰色辅助文字
TEXT_RED = RGBColor(0xC0, 0x00, 0x00)         # 红色强调
TEXT_BLACK = RGBColor(0x00, 0x00, 0x00)       # 纯黑

# 分隔/装饰
DIVIDER_GREEN = RGBColor(0x4C, 0xAF, 0x50)   # 绿色分隔线
DIVIDER_BLUE = RGBColor(0x2E, 0x75, 0xB6)    # 蓝色
DIVIDER_DARK = RGBColor(0x3B, 0x3B, 0x3B)    # 深灰

# Phase 色
PHASE_BLUE = RGBColor(0x2E, 0x75, 0xB6)      # 蓝色
PHASE_GREEN = RGBColor(0x4C, 0xAF, 0x50)      # 绿色
PHASE_ORANGE = RGBColor(0xE6, 0x7E, 0x22)    # 橙色
PHASE_PURPLE = RGBColor(0x8E, 0x44, 0xAD)    # 紫色

PHASE_COLORS = [PHASE_BLUE, PHASE_GREEN, PHASE_ORANGE, PHASE_PURPLE]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# ==================== 工具函数 ====================
def add_bg(slide, color=BG_WHITE):
    """添加纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=TEXT_DARK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='Microsoft YaHei'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rich_textbox(slide, left, top, width, height, lines, font_name='Microsoft YaHei'):
    """添加多段富文本框, lines = [(text, size, color, bold, align), ...]"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold, align) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(2)
    return txBox


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    """添加矩形色块"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.0)
    else:
        shape.line.fill.background()
    return shape


def add_plain_rect(slide, left, top, width, height, fill_color, line_color=None):
    """添加普通矩形"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.0)
    else:
        shape.line.fill.background()
    return shape


def add_arrow(slide, left, top, width, height, color):
    """添加向下箭头"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_right_arrow(slide, left, top, width, height, color):
    """添加向右箭头"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, Inches(left), Inches(top),
        Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_title_bar(slide, title, subtitle=None):
    """统一标题栏 - 浅色风格"""
    # 左侧绿色装饰条
    add_plain_rect(slide, 0, 0, 0.08, 1.1, DIVIDER_GREEN)
    # 标题文字
    add_textbox(slide, 0.4, 0.15, 11, 0.6, title, font_size=26,
                color=TEXT_DARK, bold=True)
    if subtitle:
        add_textbox(slide, 0.4, 0.7, 11, 0.35, subtitle, font_size=13,
                    color=TEXT_GRAY)
    # 底部分隔线
    add_plain_rect(slide, 0, 1.1, 13.333, 0.02, DIVIDER_GREEN)


def add_section_number(slide, text, left, top):
    """添加章节编号圆标"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left), Inches(top), Inches(0.5), Inches(0.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DIVIDER_GREEN
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)


# ==================== 封面页 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# 顶部装饰条
add_plain_rect(slide, 0, 0, 13.333, 0.08, DIVIDER_GREEN)

# 主标题
add_textbox(slide, 1.5, 2.0, 10, 1.2, '冷启动规则挖掘工具',
            font_size=44, color=TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 3.2, 10, 0.8, '从少量 Badcase 轨迹到可部署评估规则的全自动管线',
            font_size=22, color=DIVIDER_BLUE, alignment=PP_ALIGN.CENTER)

# 内容概要卡片
add_rect(slide, 2, 4.5, 9.3, 1.2, ACCENT_GREEN, DIVIDER_GREEN)
lines = [
    ('30~100条仅标注通过/失败的轨迹  →  自动产出评估规则', 15, TEXT_DARK, True, PP_ALIGN.CENTER),
    ('（定义 + IF-THEN + Few-Shot + Skill归因）', 12, TEXT_GRAY, False, PP_ALIGN.CENTER),
]
add_rich_textbox(slide, 2.2, 4.6, 8.9, 1.0, lines)

# 日期
add_textbox(slide, 4, 6.5, 5, 0.4, '2026.06', font_size=14,
            color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)


# ========================================================================
# 第一章：业界洞察
# ========================================================================

# ==================== 页面：章节分隔 - 业界洞察 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)

add_section_number(slide, '01', 6.0, 2.0)
add_textbox(slide, 1.5, 2.6, 10, 1.0, '业界洞察',
            font_size=40, color=TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 3.6, 10, 0.6, 'LLM Agent评估的学术前沿与工业实践',
            font_size=20, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

# 三个关键词卡片
kw_cards = [
    ('规则发现', '从数据中自动归纳评估规则', ACCENT_GREEN),
    ('弱监督', '低成本标注驱动高质量产出', ACCENT_PURPLE),
    ('可解释性', '规则可追溯、可验证、可执行', ACCENT_PINK),
]
for i, (kw, desc, bg) in enumerate(kw_cards):
    x = 1.5 + i * 3.7
    add_rect(slide, x, 4.5, 3.3, 1.2, bg, DIVIDER_GREEN)
    add_textbox(slide, x + 0.2, 4.6, 2.9, 0.5, kw, font_size=18, color=TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.2, 5.1, 2.9, 0.5, desc, font_size=12, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)


# ==================== 页面：LLM Agent评估的冷启动困境 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, 'LLM Agent 质量评估的冷启动困境')

# 左栏 - 行业背景
add_rect(slide, 0.5, 1.5, 5.8, 5.3, ACCENT_GREEN)
add_textbox(slide, 0.8, 1.6, 5.2, 0.5, '行业背景', font_size=20, color=DIVIDER_BLUE, bold=True)

bg_lines = [
    ('• LLM Agent 在金融、客服等高合规领域加速落地，质量保障是核心瓶颈', 13, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('', 6, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('• SWE-bench、AgentBench 等学术基准只验证能力上限，不关注合规缺陷', 13, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('', 6, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('• Red-Teaming / LLM-as-Judge 范式依赖人工编写评估规则', 13, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('  规则本身成为新的瓶颈', 13, TEXT_RED, False, PP_ALIGN.LEFT),
    ('', 6, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('• 学术界关注"如何评估"，工业界卡在"评估什么"', 13, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('  —— 即规则从哪来', 13, TEXT_RED, False, PP_ALIGN.LEFT),
]
add_rich_textbox(slide, 0.8, 2.2, 5.2, 4.4, bg_lines)

# 右栏 - 实际痛点
add_rect(slide, 7, 1.5, 5.8, 5.3, ACCENT_PINK, TEXT_RED)
add_textbox(slide, 7.3, 1.6, 5.2, 0.5, '实际痛点', font_size=20, color=TEXT_RED, bold=True)

pain_items = [
    ('❶ 规则空白', '新业务上线无现成评估规则，人工编写周期2-4周、覆盖不全'),
    ('❷ 人工归纳瓶颈', '专家逐条阅读badcase，效率低、主观性强、难以标准化'),
    ('❸ 规则不可操作', '"注意金额精度"≠"IF CHK010=0 THEN 金额精度处理不当"'),
    ('❹ 归因缺失', '发现问题无法定位到具体Skill，优化无从下手'),
]
y = 2.3
for title, desc in pain_items:
    add_textbox(slide, 7.4, y, 5.2, 0.35, title, font_size=15, color=TEXT_RED, bold=True)
    add_textbox(slide, 7.4, y + 0.35, 5.2, 0.6, desc, font_size=12, color=TEXT_DARK)
    y += 1.1

# 底部强调
add_plain_rect(slide, 0.5, 6.9, 12.3, 0.45, DIVIDER_GREEN)
add_textbox(slide, 0.5, 6.93, 12.3, 0.4,
            '30-100条badcase → 自动化 → 规则产出周期从2-4周缩短到2-4小时',
            font_size=16, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, alignment=PP_ALIGN.CENTER)


# ==================== 页面：学术对标 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '学术对标：相关研究如何解决"规则从哪来"的问题')

refs = [
    ('BERTopic', 'Grootendorst, 2022\nBERTopic: Neural topic modeling with a class-based TF-IDF procedure',
     '从大量文档中自动归纳主题结构（Topic Modeling）\n'
     '→ 本工具借鉴其思想：在Phase 3中用LLM自由聚类归纳错误类别\n'
     '→ 未来面对1000+数据时，可用BERTopic高效进行特征归纳',
     ACCENT_GREEN, PHASE_GREEN),
    ('AutoSEP', 'arXiv 2506.03195\nUnlabeled Data Improves Fine-Grained Image Zero-shot Classification',
     '证明无标注数据可提升细粒度分类效果\n'
     '→ 本工具验证了类似思想：仅需"通过/失败"二值标注\n'
     '  即可产出细粒度评估规则，无需逐条标注错误类别',
     ACCENT_PURPLE, PHASE_PURPLE),
    ('Snorkel', 'Ratner et al., 2018, VLDB\nSnorkel: Rapid Training Data Creation with Weak Supervision',
     '用弱监督（编写标注函数 LF）代替全量人工标注\n'
     '→ 本工具产出IF-THEN规则，本质上即Snorkel中的标注函数\n'
     '  冷启动规则 = 弱监督信号，驱动下游评估器',
     ACCENT_LAVENDER, PHASE_BLUE),
    ('Anchors', 'Ribeiro et al., 2018\nAnchors: High Precision Model-Agnostic Explanations',
     '用高精度的IF-THEN规则解释模型预测\n'
     '→ 本工具产出同样格式的IF-THEN条件，但目标不同：\n'
     '  Anchors解释"为何如此预测"，本工具定义"何时判定违规"',
     ACCENT_PINK, TEXT_RED),
]

for i, (topic, refs_str, desc, bg_color, tag_color) in enumerate(refs):
    x = 0.4 + (i % 2) * 6.4
    y = 1.35 + (i // 2) * 3.0
    add_rect(slide, x, y, 6.1, 2.7, bg_color, tag_color)
    # 标签
    tag_shape = add_rect(slide, x, y, 1.4, 0.42, tag_color)
    add_textbox(slide, x + 0.05, y + 0.02, 1.3, 0.38, topic, font_size=13,
                color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, alignment=PP_ALIGN.CENTER)
    # 引用
    add_textbox(slide, x + 1.5, y + 0.05, 4.4, 0.35, refs_str.split('\n')[0], font_size=9, color=TEXT_GRAY)
    # 描述
    add_textbox(slide, x + 0.15, y + 0.5, 5.8, 2.1, desc, font_size=11, color=TEXT_DARK)

# 底部关键总结
add_plain_rect(slide, 0.4, 7.05, 12.5, 0.35, DIVIDER_BLUE)
add_textbox(slide, 0.5, 7.07, 12.3, 0.3,
            '核心差异: 现有工作关注"如何评估"(评估方法), 本工具关注"评估什么"(评估规则内容) → 冷启动解法',
            font_size=13, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, alignment=PP_ALIGN.CENTER)


# ========================================================================
# 第二章：工具介绍
# ========================================================================

# ==================== 页面：章节分隔 - 工具介绍 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)

add_section_number(slide, '02', 6.0, 2.0)
add_textbox(slide, 1.5, 2.6, 10, 1.0, '工具介绍',
            font_size=40, color=TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 3.6, 10, 0.6, '冷启动规则挖掘工具的设计与实现',
            font_size=20, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

kw_cards2 = [
    ('四阶段管线', 'Phase2→3→4→5 逐层蒸馏', ACCENT_GREEN),
    ('经验元体系', '规则元·反思元·案例元', ACCENT_PURPLE),
    ('工程健壮', '断点续跑·重试降级·零丢失', ACCENT_PINK),
]
for i, (kw, desc, bg) in enumerate(kw_cards2):
    x = 1.5 + i * 3.7
    add_rect(slide, x, 4.5, 3.3, 1.2, bg, DIVIDER_GREEN)
    add_textbox(slide, x + 0.2, 4.6, 2.9, 0.5, kw, font_size=18, color=TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.2, 5.1, 2.9, 0.5, desc, font_size=12, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)


# ==================== 页面：工具定位与核心价值 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '冷启动规则挖掘工具 — 定位')

# 核心定位
add_rect(slide, 0.5, 1.4, 12.3, 0.9, ACCENT_GREEN, DIVIDER_GREEN)
add_textbox(slide, 0.7, 1.45, 11.9, 0.8,
            '从少量仅标注"通过/失败"的轨迹数据，全自动产出可部署的评估规则\n'
            '（含类别定义 + IF-THEN条件 + Few-Shot示例 + Skill归因）',
            font_size=16, color=TEXT_DARK, alignment=PP_ALIGN.CENTER)

# 闭环流程
boxes = [
    ('少量Badcase\n轨迹(30-100条)', BG_WHITE, TEXT_DARK),
    ('冷启动规则\n挖掘工具\nPhase2→3→4→5', ACCENT_GREEN, RGBColor(0xFF, 0xFF, 0xFF)),
    ('评估规则\n(JSON)', BG_WHITE, TEXT_DARK),
    ('评估器', BG_WHITE, TEXT_DARK),
    ('优化器', BG_WHITE, TEXT_DARK),
]
box_labels = ['仅标注\n通过/失败', '本工具', 'IF-THEN\nFew-Shot\nSkill归因', '通过/部分通过/失败\n+错误模式', '针对性修改\nSkill Prompt']
x_positions = [0.6, 3.0, 5.6, 8.2, 10.6]

for i, ((text, bg, tc), label) in enumerate(zip(boxes, box_labels)):
    border = DIVIDER_GREEN if i == 1 else DIVIDER_BLUE
    add_rect(slide, x_positions[i], 2.7, 2.2, 1.5, bg, border)
    add_textbox(slide, x_positions[i] + 0.05, 2.8, 2.1, 0.8, text,
                font_size=12, color=tc if bg != BG_WHITE else TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x_positions[i] + 0.05, 3.6, 2.1, 0.5, label,
                font_size=9, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
    if i < len(boxes) - 1:
        add_right_arrow(slide, x_positions[i] + 2.25, 3.2, 0.55, 0.35,
                        PHASE_COLORS[i] if i < 4 else TEXT_GRAY)

add_textbox(slide, 0.5, 4.6, 12.3, 0.4, '产出规则示例 →',
            font_size=13, color=TEXT_GRAY)

# 规则JSON示例
add_rect(slide, 0.5, 5.0, 12.3, 2.2, BG_LIGHT)
json_text = (
    '{\n'
    '  "id": "R005",\n'
    '  "error_category": "CAT005-金额精度处理不当",\n'
    '  "if_conditions": [{"feature":"CHK011_final","op":"==","value":0}, ...],\n'
    '  "skill_attribution": {"top3": [{"skill_name":"fund_planning_skill",'
    '"problematic_rule":"未规定金额精度校验","confidence":0.8}, ...]},\n'
    '  "few_shots": {"positive_examples": [...], "negative_examples": [...]},\n'
    '  "confidence": 0.85\n'
    '}'
)
add_textbox(slide, 0.7, 5.1, 10, 2.0, json_text, font_size=11, color=PHASE_BLUE)


# ==================== 页面：经验元 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '经验元：智能体执行任务的最小可复用经验单元')

# 核心定义
add_rect(slide, 0.5, 1.4, 12.3, 1.0, ACCENT_LAVENDER, PHASE_PURPLE)
add_textbox(slide, 0.7, 1.45, 11.9, 0.9,
            '经验元 = 智能体在执行任务过程中产生的"最小可复用经验单元"\n'
            '是评估规则、反思记录、黄金案例的统一抽象，支撑评估体系的持续进化',
            font_size=15, color=TEXT_DARK, alignment=PP_ALIGN.CENTER)

# 三类经验元卡片
exp_cards = [
    ('规则元', PHASE_BLUE, ACCENT_GREEN,
     ['由冷启动规则挖掘工具自动产出',
      '格式: IF-THEN条件 + Few-Shot + Skill归因',
      '来源: 30-100条标注数据→四阶段管线',
      '作用: 评估器直接执行，判定通过/失败',
      '🔄 迭代: 积累新数据后可使用规则挖掘工具刷新']),
    ('反思元', PHASE_PURPLE, ACCENT_PURPLE,
     ['由评估器反思产生',
      '格式: 错误模式描述 + 根因分析 + 改进建议',
      '来源: 评估器在评估过程中的meta认知',
      '作用: 补充规则元未能覆盖的边界case',
      '🔄 迭代: 周期性归纳→沉淀为新的规则元']),
    ('案例元', PHASE_ORANGE, ACCENT_PINK,
     ['由Few-Shot总结成为黄金语料',
      '格式: 标准化的正例/负例对话轨迹',
      '来源: 人工精选或从规则挖掘中提取典型轨迹',
      '作用: 作为评估器判定时的参考标准',
      '🔄 迭代: 评估结果反馈→更新黄金语料库']),
]

for i, (title, tag_color, bg_color, bullets) in enumerate(exp_cards):
    x = 0.5 + i * 4.2
    add_rect(slide, x, 2.7, 3.9, 4.2, bg_color, tag_color)
    # 标签
    tag_shape = add_rect(slide, x, 2.7, 1.6, 0.45, tag_color)
    add_textbox(slide, x + 0.05, 2.72, 1.5, 0.4, title, font_size=16,
                color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, alignment=PP_ALIGN.CENTER)
    # 要点
    y = 3.3
    for bullet in bullets:
        add_textbox(slide, x + 0.15, y, 3.6, 0.55, '• ' + bullet, font_size=11, color=TEXT_DARK)
        y += 0.6

# 底部：冷启动闭环
add_plain_rect(slide, 0.5, 7.05, 12.3, 0.35, DIVIDER_GREEN)
add_textbox(slide, 0.6, 7.07, 12, 0.3,
            '冷启动流程: 人工标注50-100条(通过/失败) → 规则挖掘 → 规则元产出 → 运行积累新数据 → 规则刷新',
            font_size=13, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, alignment=PP_ALIGN.CENTER)


# ==================== 页面：经验元-冷启动与迭代循环 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '经验元的冷启动与迭代进化')

# 阶段1：冷启动
add_rect(slide, 0.4, 1.4, 6.0, 2.8, ACCENT_GREEN, DIVIDER_GREEN)
add_textbox(slide, 0.6, 1.5, 5.6, 0.4, '阶段1：冷启动', font_size=18, color=DIVIDER_GREEN, bold=True)
cold_lines = [
    ('1. 人工标注50-100条轨迹数据（仅标注通过/失败）', 12, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('2. 使用规则挖掘工具产出规则元（IF-THEN + Few-Shot）', 12, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('3. 评估器加载规则元，开始自动化评估', 12, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('4. 评估器在评估过程中产出反思元', 12, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('5. 人工精选典型轨迹，形成案例元（黄金语料）', 12, TEXT_DARK, False, PP_ALIGN.LEFT),
]
add_rich_textbox(slide, 0.6, 2.0, 5.6, 2.0, cold_lines)

# 阶段2：迭代进化
add_rect(slide, 6.9, 1.4, 6.0, 2.8, ACCENT_PURPLE, PHASE_PURPLE)
add_textbox(slide, 7.1, 1.5, 5.6, 0.4, '阶段2：迭代进化', font_size=18, color=PHASE_PURPLE, bold=True)
iter_lines = [
    ('1. 评估器运行一段时间，积累新的轨迹数据', 12, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('2. 新数据 + 原有数据 → 再次使用规则挖掘工具', 12, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('3. 规则刷新：补充新规则 + 优化已有规则', 12, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('4. 反思元归纳 → 沉淀为新的规则元', 12, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('5. 案例元更新 → 评估器判定标准更精准', 12, TEXT_DARK, False, PP_ALIGN.LEFT),
]
add_rich_textbox(slide, 7.1, 2.0, 5.6, 2.0, iter_lines)

# 经验元流转图
add_rect(slide, 0.4, 4.5, 12.5, 2.8, BG_LIGHT)
add_textbox(slide, 0.6, 4.55, 12, 0.4, '经验元流转示意', font_size=16, color=TEXT_DARK, bold=True)

# 流程：数据→规则挖掘→规则元→评估器→反思元→归纳→规则元
flow_items = [
    ('标注数据\n(50-100条)', ACCENT_GREEN),
    ('规则挖掘\n工具', DIVIDER_GREEN),
    ('规则元', PHASE_BLUE),
    ('评估器', ACCENT_ORANGE_BG),
    ('反思元', PHASE_PURPLE),
    ('归纳\n→新规则元', ACCENT_PINK),
]
fx_positions = [0.5, 2.5, 4.5, 6.5, 8.5, 10.5]
for i, ((text, color), xp) in enumerate(zip(flow_items, fx_positions)):
    add_rect(slide, xp, 5.1, 1.8, 1.0, color)
    tc = RGBColor(0xFF, 0xFF, 0xFF) if color in [DIVIDER_GREEN, PHASE_BLUE, PHASE_PURPLE] else TEXT_DARK
    add_textbox(slide, xp + 0.05, 5.15, 1.7, 0.9, text, font_size=11, color=tc, bold=True, alignment=PP_ALIGN.CENTER)
    if i < len(flow_items) - 1:
        add_right_arrow(slide, xp + 1.85, 5.45, 0.45, 0.3, TEXT_GRAY)

# 底部标注
add_textbox(slide, 0.6, 6.3, 12, 0.5,
            '案例元（黄金语料）作为评估器判定的参考标准，与规则元、反思元协同工作',
            font_size=12, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)


# ==================== 页面：四阶段管线架构 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '四阶段管线架构', 'Phase 2 → 3 → 4 → 5 逐层蒸馏，从原始轨迹到可部署规则')

phases = [
    ('Phase 2: 链路归纳', PHASE_BLUE,
     '批次迭代归纳（非一次性全量）\n'
     'Batch1从零归纳 → Batch2审阅优化 → Batch3/4增量补充\n\n'
     '输入: Badcase轨迹(分批, batch=10)\n'
     '输出: 标准业务链路 + 12~15个缺失检查点'),
    ('Phase 3: 类别归纳', PHASE_GREEN,
     'LLM自由聚类 + 审查校验\n'
     '不预设类别框架, 从数据自然生长\n'
     '校验: 可判定性/客观性/证据位置\n\n'
     '输入: 缺失检查点 + 轨迹\n'
     '输出: 5个类别 + 9~11个二值检查点(CHK)'),
    ('Phase 4: 特征提取', PHASE_ORANGE,
     'LLM逐条判定 + 缓存复用\n'
     '对每条轨迹×每个检查点 → 0/1/NA\n'
     '批量模式降低5倍LLM调用成本\n\n'
     '输入: 轨迹 + 检查点定义\n'
     '输出: 特征矩阵(轨迹×检查点)'),
    ('Phase 5: 规则挖掘', PHASE_PURPLE,
     '按类别聚合 + LLM规则归纳 + Skill归因\n'
     '重试3次 + 降级兜底(零规则丢失)\n'
     'skill_name白名单校验\n\n'
     '输入: 特征矩阵 + 类别 + 轨迹 + Skill\n'
     '输出: 可部署规则(含IF-THEN+Few-Shot)'),
]

y_start = 1.35
for i, (title, color, desc) in enumerate(phases):
    y = y_start + i * 1.45
    add_rect(slide, 0.4, y, 7.8, 1.3, BG_LIGHT)
    # Phase 标签
    tag_shape = add_rect(slide, 0.4, y, 2.0, 0.45, color)
    add_textbox(slide, 0.5, y + 0.02, 1.8, 0.4, title, font_size=13,
                color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, alignment=PP_ALIGN.CENTER)
    # 描述
    add_textbox(slide, 2.5, y + 0.05, 5.5, 1.2, desc, font_size=10, color=TEXT_DARK)
    # 层间箭头
    if i < 3:
        add_arrow(slide, 4.0, y + 1.3, 0.35, 0.15, TEXT_GRAY)

# 右侧关键技术标签
add_rect(slide, 8.6, 1.35, 4.2, 5.7, BG_LIGHT)
add_textbox(slide, 8.8, 1.45, 3.8, 0.4, '关键技术', font_size=18, color=TEXT_DARK, bold=True)

techs = [
    ('批次迭代归纳', '分批处理→增量合并→全局一致\n避免一次性输入过多导致注意力稀释'),
    ('LLM自由聚类', '不预设分类框架, 从数据自然生长\n类别名称由错误模式自发生成'),
    ('二值判定+NA兜底', 'True/False可客观判定\nNA排除不相关轨迹, 不污染统计'),
    ('缓存+断点续跑', 'checkpoint+轨迹内容哈希缓存\n中断重启不重复调用LLM'),
    ('重试+降级兜底', '3次重试→降级:用统计直接构造规则\n确保零规则丢失'),
    ('格式适配器', '自动修复LLM输出格式漂移\n字符串→字典, 缺失字段补全'),
    ('轨迹格式配置化', 'trajectory_config.json自定义\n不绑定特定数据格式'),
]
y = 2.0
for label, desc in techs:
    add_textbox(slide, 8.8, y, 3.8, 0.3, '▸ ' + label, font_size=11, color=DIVIDER_BLUE, bold=True)
    add_textbox(slide, 8.8, y + 0.28, 3.8, 0.45, desc, font_size=9, color=TEXT_GRAY)
    y += 0.78


# ==================== 页面：Phase 2 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, 'Phase 2：链路归纳', '从 Badcase 轨迹反推正确链路 + 缺失检查点')

# 左侧 - 迭代流程图
add_rect(slide, 0.4, 1.5, 5.5, 5.0, ACCENT_GREEN, PHASE_BLUE)
add_textbox(slide, 0.6, 1.55, 5.1, 0.4, '批次迭代归纳流程', font_size=16, color=PHASE_BLUE, bold=True)

iter_items = [
    ('Batch 1 (10条)', 'LLM: 从零归纳', '标准业务链路(4步)\n+ 初始检查点(6-8个)', PHASE_BLUE),
    ('Batch 2 (10条)', 'LLM: 审阅+优化\n→新检查点? 合并? 修正?', '更新后链路\n+ 检查点(8-10个)', PHASE_GREEN),
    ('Batch 3/4 (剩余)', 'LLM: 增量补充', '最终链路\n+ 完整检查点(12-15个)', PHASE_ORANGE),
]
y = 2.1
for title, action, output, color in iter_items:
    add_rect(slide, 0.6, y, 5.1, 1.3, BG_WHITE, color)
    add_textbox(slide, 0.7, y + 0.05, 2.2, 0.3, title, font_size=12, color=color, bold=True)
    add_textbox(slide, 0.7, y + 0.35, 2.2, 0.5, action, font_size=10, color=TEXT_DARK)
    add_textbox(slide, 3.0, y + 0.15, 2.5, 0.9, '→ ' + output, font_size=10, color=TEXT_DARK)
    if y < 4.5:
        add_arrow(slide, 2.8, y + 1.3, 0.35, 0.15, TEXT_GRAY)
    y += 1.5

# 右侧 - 关键设计
add_rect(slide, 6.3, 1.5, 6.3, 5.0, BG_LIGHT)
add_textbox(slide, 6.5, 1.55, 5.9, 0.4, '关键设计', font_size=16, color=PHASE_BLUE, bold=True)

designs = [
    ('❶ 强制约束 vs 注意事项',
     '❌ "注意金额精度"\n✅ "金额计算保留原始精度至两位小数，禁止截断/取整/舍入"\n\n'
     'Prompt约束LLM必须输出: 触发条件 + 判定标准 + 违反动作'),
    ('❷ 批次迭代（非一次性全量）',
     '每批10条, 首批归纳, 后续审阅优化\n类似MapReduce: 分批→增量合并→全局一致\n'
     '避免一次性输入过多轨迹导致LLM注意力稀释'),
    ('❸ 格式适配器',
     'LLM输出不稳定→格式漂移(字符串/字典/字段缺失)\n'
     'adapt_phase2_to_phase3() 自动修复, 确保下游可用'),
]
y = 2.1
for title, desc in designs:
    add_textbox(slide, 6.5, y, 5.9, 0.35, title, font_size=13, color=TEXT_RED, bold=True)
    add_textbox(slide, 6.5, y + 0.4, 5.9, 1.2, desc, font_size=10, color=TEXT_DARK)
    y += 1.65

# 底部
add_plain_rect(slide, 0.4, 6.8, 12.3, 0.35, DIVIDER_GREEN)
add_textbox(slide, 0.5, 6.82, 12, 0.3,
            '实际效果: 37条轨迹 → 4批迭代 → 15个缺失检查点',
            font_size=13, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, alignment=PP_ALIGN.CENTER)


# ==================== 页面：Phase 3 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, 'Phase 3：类别归纳', '从缺失检查点归纳可判定的类别体系')

# 左侧 - 原理
add_rect(slide, 0.4, 1.5, 6.0, 5.0, ACCENT_GREEN, PHASE_GREEN)
add_textbox(slide, 0.6, 1.55, 5.6, 0.4, '两阶段归纳', font_size=16, color=PHASE_GREEN, bold=True)

steps = [
    ('Step 1: LLM 自由聚类归纳',
     '不预设分类框架, 让类别名称从数据中自然生长\n'
     '输入: Phase2的缺失检查点清单 + 原始轨迹\n'
     '输出: 5个类别 + 每类1-3个二值检查点'),
    ('Step 2: LLM 审查校验',
     '检查每个检查点是否满足三要素:\n'
     '  ① 可判定性: 能否对任意轨迹给出True/False?\n'
     '  ② 客观性: 判定标准是否不含主观评价词?\n'
     '  ③ 证据位置明确性: 审查人能否在轨迹中找到证据?\n'
     '不合格的检查点被合并或重写'),
]
y = 2.1
for title, desc in steps:
    add_textbox(slide, 0.6, y, 5.6, 0.35, title, font_size=13, color=PHASE_GREEN, bold=True)
    add_textbox(slide, 0.6, y + 0.4, 5.6, 1.5, desc, font_size=11, color=TEXT_DARK)
    y += 2.3

# 右侧 - 聚类示意
add_rect(slide, 6.8, 1.5, 5.8, 5.0, BG_LIGHT)
add_textbox(slide, 7.0, 1.55, 5.4, 0.4, '聚类示例 (15个检查点 → 5个类别)', font_size=14, color=PHASE_GREEN, bold=True)

cats_example = [
    ('CAT001 信息缺失下杜撰或引导', 'CP002, CP010 → CHK001, CHK002', PHASE_BLUE),
    ('CAT002 关键操作前强制校验缺失', 'CP004, CP005, CP007, CP_008_new_2/3 → CHK003~005', PHASE_GREEN),
    ('CAT003 跨轮槽位继承与一致性缺失', 'CP_008_new_4, CP003, CP008 → CHK006~008', PHASE_ORANGE),
    ('CAT004 购买前密码/协议确认缺失', 'CP009 → CHK009', PHASE_PURPLE),
    ('CAT005 金额精度处理不当', 'CP006, CP011, CP_008_new_5, CP003_fix → CHK010, CHK011', TEXT_RED),
]
y = 2.1
for name, mapping, color in cats_example:
    add_rect(slide, 7.0, y, 5.4, 0.8, BG_WHITE, color)
    add_textbox(slide, 7.1, y + 0.03, 5.2, 0.3, name, font_size=11, color=color, bold=True)
    add_textbox(slide, 7.1, y + 0.38, 5.2, 0.35, mapping, font_size=9, color=TEXT_GRAY)
    y += 0.88

# 底部
add_plain_rect(slide, 0.4, 6.8, 12.3, 0.35, PHASE_GREEN)
add_textbox(slide, 0.5, 6.82, 12, 0.3,
            '每个binary_checkpoint包含: description + judgment_criteria + evidence_location + 正负例轨迹',
            font_size=13, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, alignment=PP_ALIGN.CENTER)


# ==================== 页面：Phase 4 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, 'Phase 4：LLM理解式特征提取', '构建对每条轨迹的 0/1/NA 判定矩阵')

# 上部 - 矩阵示意
add_rect(slide, 0.4, 1.5, 8.5, 2.8, BG_LIGHT)
add_textbox(slide, 0.6, 1.55, 8, 0.35, '特征矩阵 (轨迹 × 检查点)', font_size=16, color=PHASE_ORANGE, bold=True)

matrix_text = (
    '              CHK001  CHK002  CHK003  CHK008  CHK009  CHK010  CHK011\n'
    'auto-9d9bc9  │  0   │  0   │  1   │  1   │  0   │  NA  │  NA  │\n'
    'auto-682bec  │  0   │  NA  │  0   │  1   │  0   │  NA  │  NA  │\n'
    'auto-e99886  │  0   │  NA  │  1   │  1   │  0   │  NA  │  NA  │\n'
    'auto-dbb7a2  │  NA  │  NA  │  1   │  1   │  0   │   0  │   0  │\n'
    '...          │  ... │  ... │  ... │  ... │  ... │  ... │  ... │\n\n'
    '0 = 违反(False)    1 = 通过(True)    NA = 不适用(该检查点与此轨迹无关)'
)
add_textbox(slide, 0.7, 2.0, 7.8, 2.2, matrix_text, font_size=10, color=TEXT_DARK)

# 右上 - 判定示意
add_rect(slide, 9.2, 1.5, 3.7, 2.8, ACCENT_PINK, PHASE_ORANGE)
add_textbox(slide, 9.4, 1.55, 3.3, 0.35, 'LLM判定过程', font_size=14, color=PHASE_ORANGE, bold=True)
add_textbox(slide, 9.4, 2.0, 3.3, 2.1,
            '输入:\n'
            '  ① 检查点定义\n'
            '     (description+\n'
            '      judgment_criteria)\n'
            '  ② 轨迹对话片段\n\n'
            '输出:\n'
            '  final: 0/1/NA\n'
            '  reason: 判定依据\n'
            '  confidence: 0~1',
            font_size=11, color=TEXT_DARK)

# 下部 - 关键设计
add_rect(slide, 0.4, 4.6, 12.5, 2.5, BG_LIGHT)
add_textbox(slide, 0.6, 4.65, 12, 0.35, '关键设计', font_size=16, color=PHASE_ORANGE, bold=True)

p4_designs = [
    ('三值判定(0/1/NA)', 'NA表示检查点与此轨迹无关(如:轨迹未涉及金额变更), 后续统计中排除, 不污染fail率'),
    ('批量判定模式', '一次LLM调用同时判定多个检查点(降低成本), --no-batch切换逐个判定(精度优先)'),
    ('缓存机制', '相同checkpoint+轨迹→不重复调用; checkpoint定义变更→自动失效; 中断重启可复用'),
    ('轨迹格式配置化', 'trajectory_config.json: 字段名/角色映射/证据规则自定义, 不绑定特定数据格式'),
]
y = 5.15
for i, (title, desc) in enumerate(p4_designs):
    x = 0.6 + (i % 2) * 6.2
    yy = y + (i // 2) * 1.0
    add_textbox(slide, x, yy, 5.8, 0.25, '▸ ' + title, font_size=12, color=PHASE_ORANGE, bold=True)
    add_textbox(slide, x, yy + 0.28, 5.8, 0.6, desc, font_size=10, color=TEXT_DARK)


# ==================== 页面：Phase 5 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, 'Phase 5：规则挖掘 + Skill归因', '从特征矩阵到可部署评估规则')

# 左侧 - 流程
add_rect(slide, 0.4, 1.5, 5.8, 5.0, ACCENT_PURPLE, PHASE_PURPLE)
add_textbox(slide, 0.6, 1.55, 5.4, 0.4, '规则挖掘流程', font_size=16, color=PHASE_PURPLE, bold=True)

flow = [
    ('按类别聚合特征矩阵', 'CAT001: CHK001 fail=86%, CHK002 fail=80%\nCAT005: CHK010 fail=100%, CHK011 fail=77%'),
    ('筛选top-K检查点+典型轨迹', 'fail率最高的K个检查点\n+ 违反分数最高的5条轨迹'),
    ('LLM规则归纳', '输入: 类别定义 + 检查点统计 + 典型轨迹 + Skill摘要\n输出: error_reason + if_conditions + skill_attribution + few_shots'),
    ('后处理', 'skill_name白名单校验 + 模糊匹配\nif_conditions补充checkpoint描述\n重试3次 + 降级兜底(零规则丢失)'),
]
y = 2.1
for i, (title, desc) in enumerate(flow):
    add_rect(slide, 0.6, y, 5.4, 1.0, BG_WHITE, PHASE_PURPLE)
    add_textbox(slide, 0.7, y + 0.03, 5.2, 0.3, f'{i+1}. {title}', font_size=12, color=PHASE_PURPLE, bold=True)
    add_textbox(slide, 0.7, y + 0.33, 5.2, 0.6, desc, font_size=10, color=TEXT_DARK)
    if i < 3:
        add_arrow(slide, 3.0, y + 1.0, 0.3, 0.12, TEXT_GRAY)
    y += 1.2

# 右侧 - 产出示例
add_rect(slide, 6.6, 1.5, 6.1, 5.0, BG_LIGHT)
add_textbox(slide, 6.8, 1.55, 5.7, 0.4, '产出规则示例 (JSON)', font_size=14, color=PHASE_PURPLE, bold=True)

rule_json = (
    '{\n'
    '  "id": "R005",\n'
    '  "error_category": "CAT005-金额精度处理不当",\n'
    '  "error_reason": "Agent对金额未做精度校验...",\n'
    '  "if_conditions": [\n'
    '    {"feature":"CHK011_final","op":"==","value":0},\n'
    '    {"feature":"CHK010_final","op":"==","value":0}\n'
    '  ],\n'
    '  "skill_attribution": {\n'
    '    "top3": [\n'
    '      {"skill_name":"fund_planning_skill",\n'
    '       "problematic_rule":"未规定金额精度校验",\n'
    '       "confidence":0.8},\n'
    '      {"skill_name":"product_select_skill",\n'
    '       "problematic_rule":"未对输入金额精度限制",\n'
    '       "confidence":0.7}\n'
    '    ]\n'
    '  },\n'
    '  "few_shots": { ...\n'
    '  },\n'
    '  "confidence": 0.85\n'
    '}'
)
add_textbox(slide, 6.8, 2.1, 5.7, 4.2, rule_json, font_size=10, color=DIVIDER_BLUE)

# 底部
add_plain_rect(slide, 0.4, 6.75, 12.3, 0.35, PHASE_PURPLE)
add_textbox(slide, 0.5, 6.77, 12, 0.3,
            'CHK010_final==0 意味着"CHK010违反时触发" → 评估器可直接执行  |  skill_name精确到文件名 → 优化器可直接定位修改  |  重试+降级 → 零规则丢失',
            font_size=12, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, alignment=PP_ALIGN.CENTER)


# ==================== 页面：工程特性 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '工程特性 — 面向生产环境的设计')

cards = [
    ('🔄', '断点续跑', 'Phase 4-5均有缓存机制\nLLM调用结果按(checkpoint+轨迹)哈希缓存\n中断重启不重复调用', ACCENT_GREEN, PHASE_BLUE),
    ('🔁', '重试+降级', 'Phase5规则挖掘支持3次重试\n全部失败自动降级: 用统计直接构造规则\n置信度0.5标注"需人工审核"', ACCENT_PURPLE, PHASE_PURPLE),
    ('📐', '格式适配', 'Phase2→3适配器\n修复LLM输出格式漂移\n字符串→字典, 缺失字段补全', ACCENT_PINK, PHASE_ORANGE),
    ('⚙️', '配置化', 'Phase4轨迹格式可自定义\ntrajectory_config.json\n字段名/角色映射/证据规则', ACCENT_LAVENDER, DIVIDER_DARK),
]

for i, (icon, title, desc, bg_color, border_color) in enumerate(cards):
    x = 0.5 + (i % 2) * 6.3
    y = 1.5 + (i // 2) * 2.8
    add_rect(slide, x, y, 5.9, 2.5, bg_color, border_color)
    add_textbox(slide, x + 0.2, y + 0.1, 5.5, 0.5, f'{icon}  {title}', font_size=18, color=border_color, bold=True)
    add_textbox(slide, x + 0.2, y + 0.7, 5.5, 1.6, desc, font_size=13, color=TEXT_DARK)

# 底部
add_rect(slide, 0.5, 6.6, 12.3, 0.6, BG_LIGHT)
add_textbox(slide, 0.7, 6.65, 11.9, 0.5,
            '典型案例(37条银行理财badcase)端到端效果:\n'
            'Phase2→15个检查点  |  Phase3→5个类别/100%覆盖  |  Phase4→407个判定  |  Phase5→5条可部署规则/零丢失',
            font_size=12, color=DIVIDER_BLUE, bold=True, alignment=PP_ALIGN.CENTER)


# ========================================================================
# 第三章：案例说明
# ========================================================================

# ==================== 页面：章节分隔 - 案例说明 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)

add_section_number(slide, '03', 6.0, 2.0)
add_textbox(slide, 1.5, 2.6, 10, 1.0, '案例说明',
            font_size=40, color=TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 3.6, 10, 0.6, '银行理财购买Agent的Badcase规则挖掘实战',
            font_size=20, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)


# ==================== 页面：案例展示 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '案例：银行理财购买 Agent 的 Badcase 规则挖掘')

# 场景
add_rect(slide, 0.4, 1.5, 5.8, 2.2, ACCENT_GREEN, DIVIDER_GREEN)
add_textbox(slide, 0.6, 1.55, 5.4, 0.4, '场景与输入', font_size=16, color=DIVIDER_GREEN, bold=True)
add_textbox(slide, 0.6, 2.0, 5.4, 1.6,
            '• 业务场景: 工行理财购买对话Agent\n'
            '• 系统组成: AgentRule + product_select_skill + fund_planning_skill\n'
            '• 输入: 37条badcase轨迹, 仅标注"通过/失败"\n'
            '• 标注成本: 人工仅需判断合格/不合格, 无需分析具体问题',
            font_size=12, color=TEXT_DARK)

# 产出概览
add_rect(slide, 6.6, 1.5, 6.1, 2.2, ACCENT_PURPLE, PHASE_PURPLE)
add_textbox(slide, 6.8, 1.55, 5.7, 0.4, '产出概览', font_size=16, color=PHASE_PURPLE, bold=True)
add_textbox(slide, 6.8, 2.0, 5.7, 1.6,
            '• Phase2: 15个缺失检查点\n'
            '• Phase3: 5个类别, 11个二值检查点, 100%轨迹覆盖\n'
            '• Phase4: 37×11=407个判定\n'
            '• Phase5: 5条可部署规则(含IF-THEN+Few-Shot+Skill归因)',
            font_size=12, color=TEXT_DARK)

# 5条规则表
add_rect(slide, 0.4, 4.0, 12.3, 2.7, BG_LIGHT)
add_textbox(slide, 0.6, 4.05, 11.9, 0.4, '5条规则一览', font_size=16, color=TEXT_DARK, bold=True)

rules_data = [
    ('R001', '信息缺失下杜撰或引导', 'CHK001=0 AND CHK002=0', 'AgentRule', '0.9'),
    ('R002', '关键操作前强制校验缺失', 'CHK003=0', 'AgentRule', '0.9'),
    ('R003', '跨轮槽位继承与一致性缺失', 'CHK008=0', 'product_recommend', '0.9'),
    ('R004', '购买前密码/协议确认缺失', 'CHK009=0', 'AgentRule', '0.9'),
    ('R005', '金额精度处理不当', 'CHK010=0 AND CHK011=0', 'fund_planning_skill', '0.8'),
]
# 表头
add_textbox(slide, 0.6, 4.5, 1.0, 0.3, '规则', font_size=11, color=DIVIDER_BLUE, bold=True)
add_textbox(slide, 1.8, 4.5, 3.0, 0.3, '类别', font_size=11, color=DIVIDER_BLUE, bold=True)
add_textbox(slide, 5.0, 4.5, 3.5, 0.3, '触发条件', font_size=11, color=DIVIDER_BLUE, bold=True)
add_textbox(slide, 8.8, 4.5, 2.5, 0.3, '首位归因Skill', font_size=11, color=DIVIDER_BLUE, bold=True)
add_textbox(slide, 11.5, 4.5, 1.0, 0.3, '置信度', font_size=11, color=DIVIDER_BLUE, bold=True)

y = 4.85
for rid, cat, cond, skill, conf in rules_data:
    add_textbox(slide, 0.6, y, 1.0, 0.3, rid, font_size=11, color=TEXT_DARK)
    add_textbox(slide, 1.8, y, 3.0, 0.3, cat, font_size=11, color=PHASE_PURPLE)
    add_textbox(slide, 5.0, y, 3.5, 0.3, cond, font_size=10, color=PHASE_BLUE)
    add_textbox(slide, 8.8, y, 2.5, 0.3, skill, font_size=11, color=PHASE_ORANGE, bold=True)
    add_textbox(slide, 11.5, y, 1.0, 0.3, conf, font_size=11, color=TEXT_DARK, alignment=PP_ALIGN.CENTER)
    y += 0.35

# 亮点
add_rect(slide, 0.4, 6.8, 12.3, 0.5, ACCENT_ORANGE_BG)
add_textbox(slide, 0.6, 6.85, 11.9, 0.4,
            '🌟 R003"跨轮槽位继承缺失"是工具自动发现的隐蔽错误模式  |  🎯 Skill归因精确到文件名, 优化器可直接定位修改',
            font_size=12, color=TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)


# ==================== 页面：规则产出详解 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '规则产出详解', '以"金额精度处理不当"为例，追踪 Phase 2→5 全链路')

steps_detail = [
    ('Phase 2\n识别检查点', PHASE_BLUE,
     'CP_008_new_5: 校验金额小数位数≤2位\nCP003_fix: 金额保留原始精度, 禁止截断'),
    ('Phase 3\n归纳类别', PHASE_GREEN,
     'CAT005: 金额精度处理不当\nCHK010: 输入精度是否超过2位小数\nCHK011: 资金筹划是否保留原始精度'),
    ('Phase 4\n特征提取', PHASE_ORANGE,
     'auto-dbb7a2: CHK010=0, CHK011=0\n  (输入5000.1234, 未拦截+转账取整)\nauto-621981: CHK010=0, CHK011=0\nauto-705be3: CHK010=NA, CHK011=NA'),
    ('Phase 5\n规则产出', PHASE_PURPLE,
     'IF CHK011_final=0 AND CHK010_final=0\nTHEN 金额精度处理不当\n归因: fund_planning_skill(0.8)'),
]

x = 0.3
for i, (label, color, desc) in enumerate(steps_detail):
    add_rect(slide, x, 1.5, 3.05, 4.0, BG_WHITE, color)
    tag = add_rect(slide, x, 1.5, 3.05, 1.0, color)
    add_textbox(slide, x + 0.1, 1.55, 2.85, 0.9, label, font_size=15,
                color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.1, 2.7, 2.85, 2.7, desc, font_size=11, color=TEXT_DARK)
    if i < 3:
        add_textbox(slide, x + 3.05, 3.0, 0.3, 0.5, '→', font_size=24, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)
    x += 3.25

# 闭环应用
add_rect(slide, 0.4, 5.8, 12.3, 1.4, BG_LIGHT)
add_textbox(slide, 0.6, 5.85, 11.9, 0.35, '闭环应用', font_size=14, color=DIVIDER_BLUE, bold=True)
add_textbox(slide, 0.6, 6.25, 11.9, 0.8,
            '评估器收到此规则 → 对新轨迹判定CHK010/CHK011 → 输出"失败: 金额精度处理不当" → 定位到fund_planning_skill → 优化器在该Skill中增加"金额精度校验≤2位小数"约束',
            font_size=12, color=TEXT_DARK)


# ========================================================================
# 第四章：展望
# ========================================================================

# ==================== 页面：章节分隔 - 展望 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)

add_section_number(slide, '04', 6.0, 2.0)
add_textbox(slide, 1.5, 2.6, 10, 1.0, '展望',
            font_size=40, color=TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 3.6, 10, 0.6, '从当前能力到规模化应用的进化路径',
            font_size=20, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)


# ==================== 页面：展望 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '总结与未来方向')

# 核心贡献
add_rect(slide, 0.5, 1.5, 6.0, 5.0, ACCENT_GREEN, DIVIDER_GREEN)
add_textbox(slide, 0.7, 1.55, 5.6, 0.4, '核心贡献', font_size=18, color=DIVIDER_GREEN, bold=True)

contribs = [
    ('❶ 全自动冷启动', '仅需"通过/失败"标注, 无需人工逐条标注错误类别'),
    ('❷ 可执行规则产出', 'IF-THEN条件 + Few-Shot, 评估器可直接使用'),
    ('❸ Skill精准归因', '定位到具体Skill文件和缺失约束, 优化器可直接修改'),
    ('❹ 经验元体系', '规则元·反思元·案例元三位一体，支撑评估体系持续进化'),
    ('❺ 工程健壮', '断点续跑 + 重试降级 + 格式适配, 生产环境可用'),
]
y = 2.1
for title, desc in contribs:
    add_textbox(slide, 0.7, y, 5.6, 0.3, title, font_size=14, color=DIVIDER_GREEN, bold=True)
    add_textbox(slide, 0.7, y + 0.33, 5.6, 0.5, desc, font_size=12, color=TEXT_DARK)
    y += 0.85

# 未来方向
add_rect(slide, 6.8, 1.5, 5.8, 5.0, ACCENT_PURPLE, PHASE_PURPLE)
add_textbox(slide, 7.0, 1.55, 5.4, 0.4, '未来方向', font_size=18, color=PHASE_PURPLE, bold=True)

futures = [
    ('🔮 海量数据高效归纳', '面对1000+数据，使用BERTopic等工具高效进行特征归纳\n'
     '替代当前LLM自由聚类，降本增效'),
    ('📊 主动学习闭环', '评估器的新判定结果反哺Phase 2-5, 迭代优化规则\n'
     '积累数据 → 规则刷新 → 评估更精准'),
    ('🌐 跨领域迁移', '同一工具框架是否适用于客服、医疗等其他对话Agent场景\n'
     '经验元体系可否跨领域复用'),
    ('📉 弱化标注依赖', '从"通过/失败"标注进一步弱化到"仅提供对话日志"\n'
     '自动发现badcase，进一步降低冷启动门槛'),
]
y = 2.1
for title, desc in futures:
    add_textbox(slide, 7.0, y, 5.4, 0.3, title, font_size=14, color=PHASE_PURPLE, bold=True)
    add_textbox(slide, 7.0, y + 0.35, 5.4, 0.8, desc, font_size=11, color=TEXT_DARK)
    y += 1.1

# 底部标语
add_plain_rect(slide, 1.5, 6.7, 10.3, 0.5, DIVIDER_GREEN)
add_textbox(slide, 1.5, 6.72, 10.3, 0.45,
            '从 30 条 badcase 到可部署评估规则，2 小时替代 2 周人工归纳',
            font_size=18, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, alignment=PP_ALIGN.CENTER)


# ==================== 保存 ====================
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), '冷启动规则挖掘工具v3.pptx')
prs.save(output_path)
print(f'PPT已保存: {output_path}')
