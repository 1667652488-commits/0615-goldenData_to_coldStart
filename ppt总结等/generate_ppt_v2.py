#!/usr/bin/env python3
"""生成 冷启动规则挖掘工具 PPT — 白色商务风格，匹配参考PPT配色"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ==================== 配色体系（匹配参考PPT）====================
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)        # 白色背景
TEXT_DARK = RGBColor(0x19, 0x1B, 0x1F)       # 正文深色 #191B1F
TEXT_GRAY = RGBColor(0x67, 0x67, 0x68)       # 灰色辅助文字 #676768
TEXT_LIGHT = RGBColor(0x99, 0x99, 0x99)       # 浅灰文字
CARD_BLUE = RGBColor(0xEB, 0xEA, 0xFD)       # 淡蓝紫卡片 #EBEAFD
CARD_GREEN = RGBColor(0xE5, 0xF2, 0xEC)      # 淡绿卡片 #E5F2EC
CARD_PINK = RGBColor(0xFC, 0xF0, 0xF4)       # 淡粉卡片 #FCF0F4
CARD_YELLOW = RGBColor(0xF8, 0xF9, 0xF3)     # 淡黄卡片 #F8F9F3
CARD_PURPLE = RGBColor(0xF1, 0xF1, 0xFD)     # 淡紫卡片 #F1F1FD
ACCENT_RED = RGBColor(0xC0, 0x00, 0x00)       # 强调红 #C00000
ACCENT_GOLD = RGBColor(0xFF, 0xC0, 0x00)     # 金色 #FFC000
ACCENT_GREEN_B = RGBColor(0x92, 0xD0, 0x50)  # 亮绿 #92D050
ACCENT_BLUE = RGBColor(0x44, 0x44, 0xC0)     # 蓝紫标题色
DIVIDER = RGBColor(0xEC, 0xEC, 0xE5)         # 分割线 #ECECE5
CARD_BORDER = RGBColor(0xDD, 0xDD, 0xDD)     # 卡片边框

# Phase专用色
PHASE_BLUE = RGBColor(0x4F, 0x80, 0xC0)
PHASE_GREEN = RGBColor(0x50, 0xA0, 0x70)
PHASE_ORANGE = RGBColor(0xD0, 0x80, 0x30)
PHASE_PURPLE = RGBColor(0x80, 0x60, 0xC0)
PHASE_COLORS = [PHASE_BLUE, PHASE_GREEN, PHASE_ORANGE, PHASE_PURPLE]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# ==================== 工具函数 ====================
def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=TEXT_DARK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rich_textbox(slide, left, top, width, height, lines, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold, align) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
    return txBox


def add_rect(slide, left, top, width, height, fill_color, line_color=None, corner=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_arrow_down(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_title(slide, title, subtitle=None):
    """统一标题: 顶部色条 + 标题文字"""
    add_rect(slide, 0, 0, 13.333, 0.08, ACCENT_RED, corner=False)
    add_textbox(slide, 0.6, 0.3, 11, 0.6, title, font_size=26, color=TEXT_DARK, bold=True)
    if subtitle:
        add_textbox(slide, 0.6, 0.85, 11, 0.35, subtitle, font_size=13, color=TEXT_GRAY)
    add_rect(slide, 0.6, 1.2, 3, 0.03, ACCENT_RED, corner=False)


def add_section_divider(slide, section_num, section_title):
    """章节分隔页"""
    add_rect(slide, 0, 0, 13.333, 7.5, CARD_BLUE, corner=False)
    add_textbox(slide, 1, 2.5, 11, 1.2, f'{section_num}', font_size=60,
                color=ACCENT_RED, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1, 3.8, 11, 1.0, section_title, font_size=36,
                color=TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)


def add_card(slide, left, top, width, height, fill_color, border_color=None):
    return add_rect(slide, left, top, width, height, fill_color, border_color or CARD_BORDER)


# ==================== 页面1：封面 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_rect(slide, 0, 0, 13.333, 0.08, ACCENT_RED, corner=False)

add_textbox(slide, 1, 2.0, 11.3, 1.0, '冷启动规则挖掘工具', font_size=40,
            color=TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1, 3.1, 11.3, 0.7, '从少量 Badcase 轨迹到可部署评估规则的全自动管线', font_size=20,
            color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

add_card(slide, 2, 4.5, 9.3, 0.7, CARD_BLUE)
add_textbox(slide, 2.2, 4.55, 8.9, 0.6,
            '30~100条仅标注通过/失败的轨迹  →  自动产出评估规则（定义 + IF-THEN + Few-Shot + Skill归因）',
            font_size=14, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 4, 6.2, 5, 0.4, '2026.06', font_size=14,
            color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)


# ==================== 页面2：章节页 — 业界洞察 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide, '01', '业界洞察')


# ==================== 页面3：为什么需要冷启动规则 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, '为什么需要冷启动规则挖掘？')

# 左栏 - 痛点
add_card(slide, 0.5, 1.5, 6.0, 5.5, CARD_PINK)
add_textbox(slide, 0.7, 1.6, 5.6, 0.4, '当前痛点', font_size=18, color=ACCENT_RED, bold=True)

pains = [
    ('❶ 规则空白', '新业务上线无现成评估规则，人工编写周期2-4周、覆盖不全'),
    ('❷ 人工归纳瓶颈', '专家逐条阅读badcase轨迹，效率低、主观性强、难以标准化'),
    ('❸ 规则不可操作', '"注意金额精度" ≠ "IF CHK010=0 THEN 金额精度处理不当"'),
    ('❹ 归因缺失', '发现问题无法定位到具体Skill，优化无从下手'),
]
y = 2.2
for title, desc in pains:
    add_textbox(slide, 0.7, y, 5.6, 0.3, title, font_size=14, color=ACCENT_RED, bold=True)
    add_textbox(slide, 0.7, y + 0.35, 5.6, 0.6, desc, font_size=12, color=TEXT_GRAY)
    y += 1.0

# 右栏 - LLM Agent评估体系
add_card(slide, 6.8, 1.5, 6.0, 5.5, CARD_BLUE)
add_textbox(slide, 7.0, 1.6, 5.6, 0.4, 'LLM Agent 评估体系的核心问题', font_size=18, color=ACCENT_BLUE, bold=True)

add_textbox(slide, 7.0, 2.2, 5.6, 3.0,
            'Anthropic 提出了 Eval-Driven Development 范式：\n\n'
            '  "Evals are the new PRDs"\n'
            '  —— 评估即新时代的产品需求文档\n\n'
            '评估体系四要素：\n'
            '  ① Task 定义（评估什么）\n'
            '  ② Trial 执行（在沙盒中运行Agent）\n'
            '  ③ Transcript 轨迹（记录思考/工具调用/观察）\n'
            '  ④ Grader 评分（Code-based / Model-based）\n\n'
            '但是——Grader的评分标准（Rubric）从哪来？\n'
            '→ 这就是"冷启动"问题',
            font_size=12, color=TEXT_GRAY)

# 底部强调
add_rect(slide, 0.5, 7.0, 12.3, 0.35, ACCENT_RED, corner=False)
add_textbox(slide, 0.6, 7.02, 12, 0.3,
            '现有工作关注"如何评估"（评估方法）,  本工具关注"评估什么"（评估规则内容）',
            font_size=14, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, alignment=PP_ALIGN.CENTER)


# ==================== 页面4：学术对标 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, '学术对标 — 冷启动规则挖掘的理论根基')

refs = [
    ('Snorkel', 'Ratner et al., 2018, VLDB',
     '弱监督数据编程框架：用户编写标注函数(LF)而非逐条标注\n'
     '本工具的检查点判定 = Snorkel的标注函数，但用LLM自动生成而非人工编写\n'
     '核心对齐：从弱监督信号(通过/失败)到精确标签(错误类别)',
     CARD_BLUE),
    ('BERTopic', 'Grootendorst, 2022',
     '基于BERT嵌入+类TF-IDF的主题建模：从文档中自动发现主题\n'
     '本工具的Phase3类别归纳 = BERTopic的主题聚类，但用LLM语义聚类代替嵌入距离\n'
     '核心对齐：从无标签数据中自动发现语义类别',
     CARD_GREEN),
    ('AutoSEP', 'arXiv 2506.03195',
     '利用无标签数据改善细粒度零样本分类：自动发现子类别\n'
     '本工具从"通过/失败"粗标签中挖掘细粒度错误模式(5个类别) = AutoSEP的子类别发现\n'
     '核心对齐：粗标签→细粒度分类的自动化',
     CARD_YELLOW),
    ('Anchors', 'Ribeiro et al., 2018',
     '高精度模型无关解释：找到足够分类的IF-THEN规则锚点\n'
     '本工具Phase5的if_conditions = Anchors的锚定规则，但面向Agent行为而非模型预测\n'
     '核心对齐：可解释的IF-THEN规则作为产出目标',
     CARD_PINK),
]

for i, (name, cite, desc, bg) in enumerate(refs):
    x = 0.5 + (i % 2) * 6.3
    y = 1.5 + (i // 2) * 2.8
    add_card(slide, x, y, 6.0, 2.55, bg)
    add_textbox(slide, x + 0.2, y + 0.1, 5.6, 0.35, f'{name}', font_size=16, color=TEXT_DARK, bold=True)
    add_textbox(slide, x + 0.2, y + 0.45, 5.6, 0.25, cite, font_size=9, color=TEXT_LIGHT)
    add_textbox(slide, x + 0.2, y + 0.75, 5.6, 1.7, desc, font_size=11, color=TEXT_GRAY)

# 底部总结
add_textbox(slide, 0.5, 7.05, 12.3, 0.3,
            '四篇工作的共同启发：用自动化/弱监督方法替代人工，从粗粒度信号中发现精细结构',
            font_size=13, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)


# ==================== 页面5：章节页 — 工具介绍 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide, '02', '工具介绍')


# ==================== 页面6：架构总览 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, '四阶段管线架构', 'Phase 2→3→4→5 逐层蒸馏，从原始轨迹到可部署规则')

phases = [
    ('Phase 2: 链路归纳', PHASE_BLUE, CARD_BLUE,
     '批次迭代归纳（非一次性全量）\n'
     'Batch1从零归纳 → Batch2审阅优化 → Batch3/4增量补充\n\n'
     '输入: Badcase轨迹(分批)\n'
     '输出: 标准链路 + 12~15个缺失检查点'),
    ('Phase 3: 类别归纳', PHASE_GREEN, CARD_GREEN,
     'LLM自由聚类 + 审查校验\n'
     '不预设类别, 从数据自然生长\n'
     '校验: 可判定性/客观性/证据位置\n\n'
     '输入: 缺失检查点 + 轨迹\n'
     '输出: 5个类别 + 9~11个二值检查点'),
    ('Phase 4: 特征提取', PHASE_ORANGE, CARD_YELLOW,
     'LLM逐条判定 + 缓存复用\n'
     '对每条轨迹×每个检查点 → 0/1/NA\n'
     '批量模式降低5倍LLM调用成本\n\n'
     '输入: 轨迹 + 检查点定义\n'
     '输出: 特征矩阵(轨迹×检查点)'),
    ('Phase 5: 规则挖掘', PHASE_PURPLE, CARD_PINK,
     '按类别聚合 + LLM规则归纳 + Skill归因\n'
     '重试3次 + 降级兜底(零规则丢失)\n'
     'skill_name白名单校验\n\n'
     '输入: 特征矩阵 + 类别 + Skill\n'
     '输出: 可部署规则(含IF-THEN+Few-Shot)'),
]

y_start = 1.5
for i, (title, accent, bg, desc) in enumerate(phases):
    y = y_start + i * 1.4
    add_card(slide, 0.4, y, 7.8, 1.25, bg, accent)
    # 左侧色条
    add_rect(slide, 0.4, y, 0.12, 1.25, accent, corner=False)
    # 标题
    add_textbox(slide, 0.65, y + 0.05, 2.3, 0.3, title, font_size=13, color=accent, bold=True)
    add_textbox(slide, 2.8, y + 0.05, 5.2, 1.15, desc, font_size=10, color=TEXT_GRAY)
    # 箭头
    if i < 3:
        add_arrow_down(slide, 3.9, y + 1.25, 0.35, 0.15, TEXT_LIGHT)

# 右侧关键技术
add_card(slide, 8.6, 1.5, 4.2, 5.5, CARD_PURPLE)
add_textbox(slide, 8.8, 1.6, 3.8, 0.4, '关键技术', font_size=18, color=TEXT_DARK, bold=True)

techs = [
    '🔄 批次迭代归纳 — 分批→增量→全局一致',
    '🎯 LLM自由聚类 — 类别从数据自然生长',
    '✅ 二值判定+NA — 客观可执行',
    '💾 缓存+断点续跑 — 中断重启不重复',
    '🔁 重试+降级兜底 — 零规则丢失',
    '⚙️ 格式适配器 — 自动修复LLM格式漂移',
    '📐 轨迹配置化 — 不绑定数据格式',
]
y = 2.2
for t in techs:
    add_textbox(slide, 8.8, y, 3.8, 0.35, t, font_size=11, color=TEXT_GRAY)
    y += 0.55


# ==================== 页面7：Phase 2 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, 'Phase 2：链路归纳', '从 Badcase 轨迹反推正确链路 + 缺失检查点')

# 左侧 - 迭代流程
add_card(slide, 0.4, 1.5, 5.5, 5.5, CARD_BLUE, PHASE_BLUE)
add_textbox(slide, 0.6, 1.6, 5.1, 0.4, '批次迭代归纳流程', font_size=16, color=PHASE_BLUE, bold=True)

iter_items = [
    ('Batch 1 (10条)', 'LLM: 从零归纳', '标准业务链路(4步)\n+ 初始检查点(6-8个)'),
    ('Batch 2 (10条)', 'LLM: 审阅+优化\n→新检查点? 合并?', '更新后链路\n+ 检查点(8-10个)'),
    ('Batch 3/4', 'LLM: 增量补充', '最终链路\n+ 完整检查点(12-15个)'),
]
y = 2.2
for title, action, output in iter_items:
    add_card(slide, 0.6, y, 5.1, 1.3, RGBColor(0xFF, 0xFF, 0xFF), PHASE_BLUE)
    add_textbox(slide, 0.7, y + 0.05, 2.0, 0.3, title, font_size=12, color=PHASE_BLUE, bold=True)
    add_textbox(slide, 0.7, y + 0.4, 2.0, 0.5, action, font_size=10, color=TEXT_GRAY)
    add_textbox(slide, 2.8, y + 0.15, 2.7, 0.9, '→ ' + output, font_size=10, color=TEXT_DARK)
    if y < 4.5:
        add_arrow_down(slide, 2.8, y + 1.3, 0.35, 0.12, PHASE_BLUE)
    y += 1.55

# 右侧 - 关键设计
add_card(slide, 6.3, 1.5, 6.3, 5.5, RGBColor(0xFF, 0xFF, 0xFF))
add_textbox(slide, 6.5, 1.6, 5.9, 0.4, '关键设计', font_size=16, color=PHASE_BLUE, bold=True)

designs = [
    ('❶ 强制约束 vs 注意事项',
     '❌ "注意金额精度"\n✅ "金额计算保留原始精度至两位小数，禁止截断/取整/舍入"\n'
     'Prompt约束: 触发条件 + 判定标准 + 违反动作'),
    ('❷ 批次迭代（非一次性全量）',
     '每批10条, 首批归纳后审阅优化\n类似MapReduce: 分批→增量合并→全局一致\n'
     '避免一次性输入过多轨迹导致LLM注意力稀释'),
    ('❸ 格式适配器',
     'LLM输出不稳定→格式漂移(字符串/字典/字段缺失)\n'
     'adapt_phase2_to_phase3() 自动修复, 确保下游可用'),
]
y = 2.2
for title, desc in designs:
    add_textbox(slide, 6.5, y, 5.9, 0.3, title, font_size=13, color=PHASE_BLUE, bold=True)
    add_textbox(slide, 6.5, y + 0.35, 5.9, 1.3, desc, font_size=10, color=TEXT_GRAY)
    y += 1.7


# ==================== 页面8：Phase 3 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, 'Phase 3：类别归纳', '从缺失检查点归纳可判定的类别体系')

add_card(slide, 0.4, 1.5, 6.0, 5.5, CARD_GREEN, PHASE_GREEN)
add_textbox(slide, 0.6, 1.6, 5.6, 0.4, '两阶段归纳', font_size=16, color=PHASE_GREEN, bold=True)

steps_p3 = [
    ('Step 1: LLM自由聚类归纳',
     '不预设分类框架, 让类别名称从数据中自然生长\n输入: Phase2的缺失检查点 + 原始轨迹\n输出: 5个类别 + 每类1-3个二值检查点'),
    ('Step 2: LLM审查校验',
     '检查三要素:\n  ① 可判定性: 能否对任意轨迹给出True/False?\n  ② 客观性: 是否不含主观评价词?\n  ③ 证据位置明确性\n不合格的检查点被合并或重写'),
]
y = 2.2
for title, desc in steps_p3:
    add_textbox(slide, 0.6, y, 5.6, 0.35, title, font_size=13, color=PHASE_GREEN, bold=True)
    add_textbox(slide, 0.6, y + 0.4, 5.6, 1.5, desc, font_size=11, color=TEXT_GRAY)
    y += 2.2

add_card(slide, 6.8, 1.5, 5.8, 5.5, RGBColor(0xFF, 0xFF, 0xFF))
add_textbox(slide, 7.0, 1.6, 5.4, 0.4, '聚类示例 (15个检查点 → 5个类别)', font_size=14, color=PHASE_GREEN, bold=True)

cats_example = [
    ('CAT001 信息缺失下杜撰或引导', 'CP002, CP010 → CHK001, CHK002'),
    ('CAT002 关键操作前强制校验缺失', 'CP004, CP005, CP007 → CHK003~005'),
    ('CAT003 跨轮槽位继承与一致性缺失', 'CP_008_new_4 → CHK006~008'),
    ('CAT004 购买前密码/协议确认缺失', 'CP009 → CHK009'),
    ('CAT005 金额精度处理不当', 'CP006, CP011, CP_008_new_5 → CHK010, CHK011'),
]
y = 2.2
for name, mapping in cats_example:
    add_card(slide, 7.0, y, 5.4, 0.85, CARD_GREEN, PHASE_GREEN)
    add_textbox(slide, 7.1, y + 0.05, 5.2, 0.3, name, font_size=11, color=PHASE_GREEN, bold=True)
    add_textbox(slide, 7.1, y + 0.4, 5.2, 0.35, mapping, font_size=9, color=TEXT_LIGHT)
    y += 1.0


# ==================== 页面9：Phase 4 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, 'Phase 4：LLM理解式特征提取', '构建对每条轨迹的 0/1/NA 判定矩阵')

add_card(slide, 0.4, 1.5, 8.5, 2.8, CARD_YELLOW, PHASE_ORANGE)
add_textbox(slide, 0.6, 1.6, 8, 0.35, '特征矩阵 (轨迹 × 检查点)', font_size=16, color=PHASE_ORANGE, bold=True)

matrix_text = (
    '              CHK001  CHK002  CHK003  CHK008  CHK009  CHK010  CHK011\n'
    'auto-9d9bc9  │  0   │  0   │  1   │  1   │  0   │  NA  │  NA  │\n'
    'auto-682bec  │  0   │  NA  │  0   │  1   │  0   │  NA  │  NA  │\n'
    'auto-e99886  │  0   │  NA  │  1   │  1   │  0   │  NA  │  NA  │\n'
    'auto-dbb7a2  │  NA  │  NA  │  1   │  1   │  0   │   0  │   0  │\n\n'
    '0 = 违反(False)    1 = 通过(True)    NA = 不适用'
)
add_textbox(slide, 0.7, 2.1, 7.8, 2.0, matrix_text, font_size=10, color=TEXT_GRAY, font_name='Consolas')

add_card(slide, 9.2, 1.5, 3.7, 2.8, RGBColor(0xFF, 0xFF, 0xFF))
add_textbox(slide, 9.4, 1.6, 3.3, 0.35, 'LLM判定过程', font_size=14, color=PHASE_ORANGE, bold=True)
add_textbox(slide, 9.4, 2.0, 3.3, 2.0,
            '输入:\n  ① 检查点定义\n     (description +\n      judgment_criteria)\n  ② 轨迹对话片段\n\n输出:\n  final: 0/1/NA\n  reason: 判定依据\n  confidence: 0~1',
            font_size=11, color=TEXT_GRAY)

add_card(slide, 0.4, 4.6, 12.5, 2.6, RGBColor(0xFF, 0xFF, 0xFF))
add_textbox(slide, 0.6, 4.7, 12, 0.35, '关键设计', font_size=16, color=PHASE_ORANGE, bold=True)

p4_items = [
    ('三值判定(0/1/NA)', 'NA表示检查点与此轨迹无关, 后续统计排除, 不污染fail率'),
    ('批量判定模式', '一次LLM调用判多个检查点降低成本, --no-batch逐个判定精度优先'),
    ('缓存机制', '相同checkpoint+轨迹→不重复调用; 定义变更→自动失效; 中断重启可复用'),
    ('轨迹格式配置化', 'trajectory_config.json: 字段名/角色映射/证据规则自定义'),
]
for i, (title, desc) in enumerate(p4_items):
    x = 0.6 + (i % 2) * 6.2
    y = 5.2 + (i // 2) * 0.95
    add_textbox(slide, x, y, 5.8, 0.25, '▸ ' + title, font_size=12, color=PHASE_ORANGE, bold=True)
    add_textbox(slide, x, y + 0.28, 5.8, 0.5, desc, font_size=10, color=TEXT_GRAY)


# ==================== 页面10：Phase 5 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, 'Phase 5：规则挖掘 + Skill归因', '从特征矩阵到可部署评估规则')

add_card(slide, 0.4, 1.5, 5.8, 5.5, CARD_PINK, PHASE_PURPLE)
add_textbox(slide, 0.6, 1.6, 5.4, 0.4, '规则挖掘流程', font_size=16, color=PHASE_PURPLE, bold=True)

flow = [
    ('1. 按类别聚合特征矩阵', 'CAT005: CHK010 fail=100%, CHK011 fail=77%'),
    ('2. 筛选top-K + 典型轨迹', 'fail率最高的K个检查点 + 违反分数最高的5条轨迹'),
    ('3. LLM规则归纳', '输入: 类别定义 + 统计 + 轨迹 + Skill摘要\n输出: error_reason + if_conditions + skill_attribution'),
    ('4. 后处理', 'skill_name白名单校验 + 模糊匹配\nif_conditions补充描述\n重试3次 + 降级兜底'),
]
y = 2.1
for title, desc in flow:
    add_card(slide, 0.6, y, 5.4, 1.0, RGBColor(0xFF, 0xFF, 0xFF), PHASE_PURPLE)
    add_textbox(slide, 0.7, y + 0.03, 5.2, 0.3, title, font_size=12, color=PHASE_PURPLE, bold=True)
    add_textbox(slide, 0.7, y + 0.38, 5.2, 0.55, desc, font_size=10, color=TEXT_GRAY)
    if y < 4.5:
        add_arrow_down(slide, 2.8, y + 1.0, 0.3, 0.12, PHASE_PURPLE)
    y += 1.2

add_card(slide, 6.6, 1.5, 6.1, 5.5, RGBColor(0xFF, 0xFF, 0xFF))
add_textbox(slide, 6.8, 1.6, 5.7, 0.4, '产出规则示例 (JSON)', font_size=14, color=PHASE_PURPLE, bold=True)

rule_json = (
    '{\n'
    '  "id": "R005",\n'
    '  "error_category": "CAT005-金额精度处理不当",\n'
    '  "if_conditions": [\n'
    '    {"feature":"CHK011_final","op":"==","value":0},\n'
    '    {"feature":"CHK010_final","op":"==","value":0}\n'
    '  ],\n'
    '  "skill_attribution": {\n'
    '    "top3": [\n'
    '      {"skill_name":"fund_planning_skill",\n'
    '       "problematic_rule":"未规定金额精度校验",\n'
    '       "confidence":0.8},\n'
    '      {"skill_name":"product_select_skill",\n'
    '       "problematic_rule":"未对输入金额精度限制",\n'
    '       "confidence":0.7}\n'
    '    ]\n'
    '  },\n'
    '  "few_shots": { ... },\n'
    '  "confidence": 0.85\n'
    '}'
)
add_textbox(slide, 6.8, 2.1, 5.7, 4.5, rule_json, font_size=10, color=PHASE_PURPLE, font_name='Consolas')

add_rect(slide, 0.4, 7.05, 12.5, 0.03, PHASE_PURPLE, corner=False)


# ==================== 页面11：经验元 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, '"经验元" — 最小可复用经验单元', '智能体在执行任务过程中产生的结构化经验沉淀')

# 核心定义
add_card(slide, 0.5, 1.5, 12.3, 1.3, CARD_BLUE, ACCENT_BLUE)
add_textbox(slide, 0.7, 1.6, 11.9, 1.1,
            '经验元是智能体在执行任务过程中产生的"最小可复用经验单元"\n'
            '可分为三类：规则元、反思元、案例元 —— 共同构成智能体的经验知识库',
            font_size=15, color=TEXT_DARK, bold=False, alignment=PP_ALIGN.CENTER)

# 三类经验元
yuan_items = [
    ('规则元', '用冷启动规则挖掘工具产生', PHASE_BLUE, CARD_BLUE,
     '• 从50-100条badcase轨迹自动挖掘\n'
     '• IF-THEN条件 + 判定标准\n'
     '• Skill归因 + Few-Shot\n'
     '• 可直接输入评估器执行'),
    ('反思元', '评估器反思产生', PHASE_GREEN, CARD_GREEN,
     '• 评估器在判定badcase时反思\n'
     '• 发现有新错误模式后提炼\n'
     '• 补充到规则库中\n'
     '• 持续丰富评估覆盖面'),
    ('案例元', 'few-shots总结为黄金语料', PHASE_ORANGE, CARD_YELLOW,
     '• 正负例对话片段\n'
     '• 标注对/错行为及原因\n'
     '• 构成黄金语料库\n'
     '• 作为规则判定的佐证参考'),
]
for i, (name, source, accent, bg, desc) in enumerate(yuan_items):
    x = 0.5 + i * 4.2
    add_card(slide, x, 3.1, 3.9, 3.2, bg, accent)
    add_rect(slide, x, 3.1, 3.9, 0.08, accent, corner=False)
    add_textbox(slide, x + 0.2, 3.3, 3.5, 0.4, name, font_size=18, color=accent, bold=True)
    add_textbox(slide, x + 0.2, 3.7, 3.5, 0.3, source, font_size=11, color=TEXT_GRAY)
    add_textbox(slide, x + 0.2, 4.1, 3.5, 2.0, desc, font_size=12, color=TEXT_GRAY)

# 冷启动循环
add_card(slide, 0.5, 6.5, 12.3, 0.7, RGBColor(0xFF, 0xFF, 0xFF))
add_textbox(slide, 0.7, 6.55, 11.9, 0.6,
            '冷启动循环: 初始50-100条人工标注(通过/失败) → 规则挖掘 → 运行积累新数据 → 规则刷新',
            font_size=14, color=ACCENT_RED, bold=True, alignment=PP_ALIGN.CENTER)


# ==================== 页面12：章节页 — 案例说明 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide, '03', '案例说明')


# ==================== 页面13：案例展示 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, '案例：银行理财购买 Agent')

add_card(slide, 0.4, 1.5, 5.8, 2.0, CARD_BLUE, ACCENT_BLUE)
add_textbox(slide, 0.6, 1.6, 5.4, 0.4, '场景与输入', font_size=16, color=ACCENT_BLUE, bold=True)
add_textbox(slide, 0.6, 2.1, 5.4, 1.3,
            '• 业务: 工行理财购买对话Agent\n'
            '• 组成: AgentRule + product_select_skill + fund_planning_skill\n'
            '• 输入: 37条badcase轨迹, 仅标注"通过/失败"\n'
            '• 标注成本: 人工仅需判断合格/不合格',
            font_size=12, color=TEXT_GRAY)

add_card(slide, 6.6, 1.5, 6.1, 2.0, CARD_GREEN, PHASE_GREEN)
add_textbox(slide, 6.8, 1.6, 5.7, 0.4, '产出概览', font_size=16, color=PHASE_GREEN, bold=True)
add_textbox(slide, 6.8, 2.1, 5.7, 1.3,
            '• Phase2: 15个缺失检查点\n'
            '• Phase3: 5个类别, 100%轨迹覆盖\n'
            '• Phase4: 37×11=407个判定\n'
            '• Phase5: 5条可部署规则',
            font_size=12, color=TEXT_GRAY)

# 规则表
add_card(slide, 0.4, 3.8, 12.3, 2.8, RGBColor(0xFF, 0xFF, 0xFF))
add_textbox(slide, 0.6, 3.9, 11.9, 0.4, '5条规则一览', font_size=16, color=TEXT_DARK, bold=True)

rules_data = [
    ('R001', '信息缺失下杜撰或引导', 'CHK001=0 AND CHK002=0', 'AgentRule', '0.9'),
    ('R002', '关键操作前强制校验缺失', 'CHK003=0', 'AgentRule', '0.9'),
    ('R003', '跨轮槽位继承与一致性缺失', 'CHK008=0', 'product_recommend', '0.9'),
    ('R004', '购买前密码/协议确认缺失', 'CHK009=0', 'AgentRule', '0.9'),
    ('R005', '金额精度处理不当', 'CHK010=0 AND CHK011=0', 'fund_planning_skill', '0.8'),
]
add_textbox(slide, 0.6, 4.35, 1.0, 0.3, '规则', font_size=11, color=ACCENT_BLUE, bold=True)
add_textbox(slide, 1.8, 4.35, 3.0, 0.3, '类别', font_size=11, color=ACCENT_BLUE, bold=True)
add_textbox(slide, 5.0, 4.35, 3.5, 0.3, '触发条件', font_size=11, color=ACCENT_BLUE, bold=True)
add_textbox(slide, 8.8, 4.35, 2.5, 0.3, '首位归因Skill', font_size=11, color=ACCENT_BLUE, bold=True)
add_textbox(slide, 11.5, 4.35, 1.0, 0.3, '置信度', font_size=11, color=ACCENT_BLUE, bold=True)

y = 4.7
for rid, cat, cond, skill, conf in rules_data:
    add_textbox(slide, 0.6, y, 1.0, 0.3, rid, font_size=11, color=TEXT_DARK)
    add_textbox(slide, 1.8, y, 3.0, 0.3, cat, font_size=11, color=PHASE_PURPLE)
    add_textbox(slide, 5.0, y, 3.5, 0.3, cond, font_size=10, color=PHASE_GREEN, font_name='Consolas')
    add_textbox(slide, 8.8, y, 2.5, 0.3, skill, font_size=11, color=PHASE_ORANGE, bold=True)
    add_textbox(slide, 11.5, y, 1.0, 0.3, conf, font_size=11, color=TEXT_DARK, alignment=PP_ALIGN.CENTER)
    y += 0.35

# 亮点
add_card(slide, 0.4, 6.8, 12.3, 0.5, CARD_PINK, ACCENT_RED)
add_textbox(slide, 0.6, 6.85, 11.9, 0.4,
            '🌟 R003"跨轮槽位继承缺失"是工具自动发现的隐蔽错误模式    |    🎯 Skill归因精确到文件名, 优化器可直接定位修改',
            font_size=12, color=ACCENT_RED, bold=True, alignment=PP_ALIGN.CENTER)


# ==================== 页面14：规则全链路追踪 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, '规则产出详解', '以"金额精度处理不当"为例，追踪 Phase 2→5 全链路')

steps_detail = [
    ('Phase 2\n识别检查点', PHASE_BLUE, CARD_BLUE,
     'CP_008_new_5:\n校验金额小数位数≤2位\n\nCP003_fix:\n金额保留原始精度,\n禁止截断'),
    ('Phase 3\n归纳类别', PHASE_GREEN, CARD_GREEN,
     'CAT005:\n金额精度处理不当\n\nCHK010: 输入精度\n是否超过2位小数\n\nCHK011: 资金筹划\n是否保留原始精度'),
    ('Phase 4\n特征提取', PHASE_ORANGE, CARD_YELLOW,
     'auto-dbb7a2:\nCHK010=0, CHK011=0\n(输入5000.1234,\n未拦截+转账取整)\n\nauto-705be3:\nCHK010=NA, CHK011=NA'),
    ('Phase 5\n规则产出', PHASE_PURPLE, CARD_PINK,
     'IF CHK011_final=0\nAND CHK010_final=0\nTHEN 金额精度处理不当\n\n归因:\nfund_planning_skill(0.8)\nproduct_select_skill(0.7)'),
]
x = 0.3
for i, (label, accent, bg, desc) in enumerate(steps_detail):
    add_card(slide, x, 1.5, 3.05, 4.3, bg, accent)
    add_rect(slide, x, 1.5, 3.05, 0.9, accent, corner=False)
    add_textbox(slide, x + 0.1, 1.55, 2.85, 0.8, label, font_size=15, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x + 0.1, 2.5, 2.85, 3.2, desc, font_size=11, color=TEXT_GRAY)
    if i < 3:
        add_textbox(slide, x + 3.05, 3.0, 0.3, 0.5, '→', font_size=22, color=TEXT_LIGHT, alignment=PP_ALIGN.CENTER)
    x += 3.25

# 闭环
add_card(slide, 0.4, 6.0, 12.3, 1.1, CARD_BLUE, ACCENT_BLUE)
add_textbox(slide, 0.6, 6.1, 11.9, 0.9,
            '闭环应用:\n'
            '评估器收到此规则 → 对新轨迹判定CHK010/CHK011 → 输出"失败: 金额精度处理不当" '
            '→ 定位到fund_planning_skill → 优化器增加"金额精度校验≤2位小数"约束',
            font_size=12, color=TEXT_DARK)


# ==================== 页面15：工程特性 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, '工程特性 — 面向生产环境的设计')

eng_cards = [
    ('🔄 断点续跑', 'Phase 4-5均有缓存\nLLM调用结果按(checkpoint+轨迹)\n哈希缓存, 中断重启不重复', PHASE_BLUE, CARD_BLUE),
    ('🔁 重试+降级', 'Phase5规则挖掘支持3次重试\n失败后自动降级: 用统计直接构造\n置信度0.5标注"需人工审核"', PHASE_GREEN, CARD_GREEN),
    ('📐 格式适配', 'Phase2→3适配器\n修复LLM输出格式漂移\n字符串→字典,缺失字段补全', PHASE_ORANGE, CARD_YELLOW),
    ('⚙️ 配置化', 'Phase4轨迹格式可自定义\ntrajectory_config.json\n字段名/角色映射/证据规则', PHASE_PURPLE, CARD_PINK),
]
for i, (title, desc, accent, bg) in enumerate(eng_cards):
    x = 0.5 + (i % 2) * 6.3
    y = 1.5 + (i // 2) * 2.7
    add_card(slide, x, y, 5.9, 2.35, bg, accent)
    add_textbox(slide, x + 0.2, y + 0.1, 5.5, 0.4, title, font_size=17, color=accent, bold=True)
    add_textbox(slide, x + 0.2, y + 0.6, 5.5, 1.6, desc, font_size=12, color=TEXT_GRAY)

add_card(slide, 0.5, 6.6, 12.3, 0.6, CARD_BLUE, ACCENT_BLUE)
add_textbox(slide, 0.7, 6.65, 11.9, 0.5,
            '端到端效果(37条银行理财badcase): 15个检查点 → 5个类别/100%覆盖 → 407个判定 → 5条可部署规则/零丢失',
            font_size=13, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)


# ==================== 页面16：展望 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, '展望 — 从冷启动到持续进化')

# 上部 - 三个演进方向
futures = [
    ('🔮 主动学习闭环', CARD_BLUE, PHASE_BLUE,
     '评估器的新判定结果反哺Phase 2-5\n迭代优化规则, 越用越精准\n经验元（规则元/反思元/案例元）\n持续丰富智能体经验知识库'),
    ('📊 海量数据高效归纳', CARD_GREEN, PHASE_GREEN,
     '面对1000+海量轨迹数据\n引入BERTopic等主题建模工具\n高效进行特征归纳与聚类\n替代纯LLM聚类, 降低成本提升速度'),
    ('📉 弱化标注依赖', CARD_PINK, PHASE_ORANGE,
     '从"通过/失败"标注\n进一步弱化到"仅提供对话日志"\n自动发现badcase, 自动挖掘规则\n实现真正的零冷启动'),
]
for i, (title, bg, accent, desc) in enumerate(futures):
    x = 0.5 + i * 4.2
    add_card(slide, x, 1.5, 3.9, 3.5, bg, accent)
    add_rect(slide, x, 1.5, 3.9, 0.08, accent, corner=False)
    add_textbox(slide, x + 0.2, 1.7, 3.5, 0.4, title, font_size=16, color=accent, bold=True)
    add_textbox(slide, x + 0.2, 2.2, 3.5, 2.5, desc, font_size=13, color=TEXT_GRAY)

# 底部 - 闭环示意
add_card(slide, 0.5, 5.3, 12.3, 1.8, CARD_BLUE, ACCENT_BLUE)
add_textbox(slide, 0.7, 5.4, 11.9, 0.4, '从冷启动到持续进化的闭环', font_size=16, color=ACCENT_BLUE, bold=True)
add_textbox(slide, 0.7, 5.9, 11.9, 1.0,
            '初始50-100条标注 → 冷启动规则挖掘 → 评估器上线 → 反思积累经验元 → 运行积累新数据(1000+) → BERTopic高效归纳 → 规则刷新 → 评估器升级 → 循环持续进化',
            font_size=13, color=TEXT_DARK)

add_rect(slide, 1.5, 7.15, 10.3, 0.04, ACCENT_RED, corner=False)
add_textbox(slide, 1.5, 7.2, 10.3, 0.25,
            '从 30 条 badcase 到可部署评估规则，2 小时替代 2 周人工归纳',
            font_size=14, color=ACCENT_RED, bold=True, alignment=PP_ALIGN.CENTER)


# ==================== 保存 ====================
output_path = os.path.join(os.path.dirname(__file__) or '.', '冷启动规则挖掘工具.pptx')
prs.save(output_path)
print(f'PPT已保存: {output_path}')
