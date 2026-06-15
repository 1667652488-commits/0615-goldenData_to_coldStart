#!/usr/bin/env python3
"""生成 冷启动规则挖掘工具 v4 PPT
聚焦：黄金数据标注 × 冷数据规则提取 联调流水线
配套提示词说明：kimi.md
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# ==================== 颜色体系 ====================
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT = RGBColor(0xF8, 0xF9, 0xF3)
BG_LIGHT_GRAY = RGBColor(0xEC, 0xEC, 0xE5)

ACCENT_GREEN = RGBColor(0xE5, 0xF2, 0xEC)
ACCENT_BLUE = RGBColor(0xEB, 0xEA, 0xFD)
ACCENT_PINK = RGBColor(0xFC, 0xF0, 0xF4)
ACCENT_YELLOW = RGBColor(0xF8, 0xF9, 0xF3)
ACCENT_PURPLE = RGBColor(0xF1, 0xF1, 0xFD)

TEXT_DARK = RGBColor(0x19, 0x1B, 0x1F)
TEXT_GRAY = RGBColor(0x67, 0x67, 0x68)
TEXT_RED = RGBColor(0xC0, 0x00, 0x00)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)

DIVIDER_GREEN = RGBColor(0x4C, 0xAF, 0x50)
DIVIDER_BLUE = RGBColor(0x2E, 0x75, 0xB6)
DIVIDER_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
DIVIDER_PURPLE = RGBColor(0x8E, 0x44, 0xAD)
DIVIDER_DARK = RGBColor(0x3B, 0x3B, 0x3B)

PHASE_BLUE = RGBColor(0x2E, 0x75, 0xB6)
PHASE_GREEN = RGBColor(0x4C, 0xAF, 0x50)
PHASE_ORANGE = RGBColor(0xE6, 0x7E, 0x22)
PHASE_PURPLE = RGBColor(0x8E, 0x44, 0xAD)
PHASE_RED = RGBColor(0xC0, 0x00, 0x00)

PHASE_COLORS = [PHASE_BLUE, PHASE_GREEN, PHASE_ORANGE, PHASE_PURPLE]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


# ==================== 工具函数 ====================
def add_bg(slide, color=BG_WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=TEXT_DARK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rich_textbox(slide, left, top, width, height, lines, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold, align) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        p.space_after = Pt(2)
    return txBox


def add_rect(slide, left, top, width, height, fill_color, line_color=None, corner=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if corner else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.0)
    else:
        shape.line.fill.background()
    return shape


def add_plain_rect(slide, left, top, width, height, fill_color, line_color=None):
    return add_rect(slide, left, top, width, height, fill_color, line_color, corner=False)


def add_arrow_right(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_arrow_down(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_title_bar(slide, title, subtitle=None):
    add_plain_rect(slide, 0, 0, 0.08, 1.1, DIVIDER_GREEN)
    add_textbox(slide, 0.4, 0.15, 11, 0.6, title, font_size=26,
                color=TEXT_DARK, bold=True)
    if subtitle:
        add_textbox(slide, 0.4, 0.7, 11, 0.35, subtitle, font_size=13,
                    color=TEXT_GRAY)
    add_plain_rect(slide, 0, 1.1, 13.333, 0.02, DIVIDER_GREEN)


def add_section_number(slide, text, left, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top),
                                    Inches(0.5), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DIVIDER_GREEN
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER


def add_section_divider(slide, num, title, subtitle):
    add_bg(slide, BG_LIGHT)
    add_section_number(slide, num, 6.0, 2.0)
    add_textbox(slide, 1.5, 2.6, 10, 1.0, title,
                font_size=40, color=TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1.5, 3.6, 10, 0.6, subtitle,
                font_size=20, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)


def add_bottom_bar(slide, text, color=DIVIDER_GREEN):
    add_plain_rect(slide, 0.5, 6.9, 12.3, 0.45, color)
    add_textbox(slide, 0.5, 6.93, 12.3, 0.4, text,
                font_size=14, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ==================== 第 1 页：封面 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_plain_rect(slide, 0, 0, 13.333, 0.08, DIVIDER_GREEN)

add_textbox(slide, 1.5, 2.0, 10, 1.2, '冷启动规则挖掘工具 v4',
            font_size=44, color=TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 3.2, 10, 0.8,
            '黄金数据标注 × 冷数据规则提取 联调流水线',
            font_size=24, color=DIVIDER_BLUE, alignment=PP_ALIGN.CENTER)

add_rect(slide, 2, 4.5, 9.3, 1.2, ACCENT_GREEN, DIVIDER_GREEN)
lines = [
    ('从少量无标注轨迹 → 可部署评估规则，端到端自动化', 15, TEXT_DARK, True, PP_ALIGN.CENTER),
    ('（黄金数据标注 + 桥接筛选 + Phase2~6 规则挖掘）', 12, TEXT_GRAY, False, PP_ALIGN.CENTER),
]
add_rich_textbox(slide, 2.2, 4.7, 8.9, 0.8, lines)

add_textbox(slide, 4, 6.5, 5, 0.4, '2026.06', font_size=14,
            color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)


# ==================== 第 2 页：目录 / 核心价值 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '核心价值与内容概览')

# 左侧核心价值卡片
values = [
    ('端到端联调', 'Phase1 黄金标注 → Phase2-6 规则挖掘\n一键打通完整流水线', ACCENT_GREEN, DIVIDER_GREEN),
    ('低成本冷启动', '仅需"通过/失败"直觉标注\n无需专家逐条归纳错误类别', ACCENT_BLUE, DIVIDER_BLUE),
    ('可执行产出', 'IF-THEN + Few-Shot + Skill 归因\n评估器可直接加载执行', ACCENT_PINK, DIVIDER_PURPLE),
    ('方法论沉淀', '系统理解 6 点 + 最终状态原则\n提示词工程可复现', ACCENT_YELLOW, DIVIDER_ORANGE),
]
for i, (title, desc, bg, border) in enumerate(values):
    x = 0.5 + (i % 2) * 6.3
    y = 1.5 + (i // 2) * 2.2
    add_rect(slide, x, y, 5.9, 1.9, bg, border)
    add_textbox(slide, x + 0.2, y + 0.15, 5.5, 0.45, title,
                font_size=16, color=border, bold=True)
    add_textbox(slide, x + 0.2, y + 0.65, 5.5, 1.1, desc,
                font_size=11, color=TEXT_DARK)

# 右侧目录时间线
add_rect(slide, 0.5, 6.2, 12.3, 1.0, BG_LIGHT)
add_textbox(slide, 0.7, 6.3, 11.9, 0.8,
            '内容结构：01 背景与痛点  →  02 完整流水线设计  →  03 方法论沉淀  →  04 实战案例  →  05 核心价值与展望',
            font_size=14, color=TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)


# ==================== 第 3 页：章节分隔 01 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide, '01', '背景与痛点', '为什么需要"标注 + 规则挖掘"联调方案？')


# ==================== 第 4 页：LLM Agent 评估的冷启动困境 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, 'LLM Agent 评估的冷启动困境')

# 左栏
add_rect(slide, 0.5, 1.5, 5.8, 5.3, ACCENT_GREEN, DIVIDER_GREEN)
add_textbox(slide, 0.8, 1.6, 5.2, 0.5, '行业背景', font_size=20, color=DIVIDER_BLUE, bold=True)
bg_lines = [
    ('• LLM Agent 在金融、客服等高合规领域加速落地，质量保障是核心瓶颈', 13, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('', 6, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('• SWE-bench、AgentBench 等学术基准只验证能力上限，不关注合规缺陷', 13, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('', 6, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('• Red-Teaming / LLM-as-Judge 范式依赖人工编写评估规则', 13, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('  规则本身成为新的瓶颈', 13, TEXT_RED, False, PP_ALIGN.LEFT),
    ('', 6, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('• 学术界关注"如何评估"，工业界卡在"评估什么"', 13, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('  —— 即规则从哪来', 13, TEXT_RED, False, PP_ALIGN.LEFT),
]
add_rich_textbox(slide, 0.8, 2.2, 5.2, 4.4, bg_lines)

# 右栏
add_rect(slide, 7, 1.5, 5.8, 5.3, ACCENT_PINK, TEXT_RED)
add_textbox(slide, 7.3, 1.6, 5.2, 0.5, '实际痛点', font_size=20, color=TEXT_RED, bold=True)
pain_items = [
    ('❶ 规则空白', '新业务上线无现成评估规则，人工编写周期 2-4 周、覆盖不全'),
    ('❷ 标注成本高', '专家逐条阅读 badcase 并标注错误类别，效率低、主观性强'),
    ('❸ 工具割裂', '黄金数据标注 与 冷数据规则提取 各自为战，数据流断裂'),
    ('❹ 规则不可操作', '"注意金额精度" ≠ "IF CHK010=0 THEN 金额精度处理不当"'),
]
y = 2.3
for title, desc in pain_items:
    add_textbox(slide, 7.4, y, 5.2, 0.35, title, font_size=15, color=TEXT_RED, bold=True)
    add_textbox(slide, 7.4, y + 0.35, 5.2, 0.6, desc, font_size=12, color=TEXT_DARK)
    y += 1.1

add_bottom_bar(slide, '30-100 条轨迹 + 2 小时 → 替代 2 周人工归纳')


# ==================== 第 5 页：两个工具的割裂与联调必要性 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '从"两张皮"到"一条线"')

# 上方表格
add_rect(slide, 0.5, 1.4, 12.3, 2.4, BG_LIGHT)
add_textbox(slide, 0.7, 1.5, 11.9, 0.4, '两个工具的现状对比', font_size=16, color=TEXT_DARK, bold=True)

# 表头
headers = ['工具', '输入', '输出', '局限']
xs = [0.7, 3.0, 6.0, 9.5]
for x, h in zip(xs, headers):
    add_textbox(slide, x, 1.95, 2.5, 0.3, h, font_size=12, color=DIVIDER_BLUE, bold=True)

# 数据行
rows = [
    ['黄金数据标注', '无标注轨迹', 'golden_output.jsonl\n(result + reason)', '只评判，不产规则'],
    ['冷数据规则提取', '失败/部分通过轨迹', 'rules.json + 自然语言规则', '依赖人工预筛选'],
]
for i, row in enumerate(rows):
    y = 2.35 + i * 0.6
    for x, cell in zip(xs, row):
        add_textbox(slide, x, y, 2.8, 0.55, cell, font_size=11, color=TEXT_DARK)

# 下方流程图
add_rect(slide, 0.5, 4.1, 12.3, 3.0, ACCENT_GREEN, DIVIDER_GREEN)
add_textbox(slide, 0.7, 4.2, 11.9, 0.4, '联调后：无标注轨迹 → 可部署规则的完整流水线', font_size=16, color=DIVIDER_GREEN, bold=True)

flow_items = [
    ('无标注\n轨迹', BG_WHITE),
    ('Phase1\n黄金数据标注', DIVIDER_BLUE),
    ('golden_\noutput.jsonl', BG_WHITE),
    ('桥接筛选\n(失败/部分通过)', DIVIDER_ORANGE),
    ('badcase_\ntraces/', BG_WHITE),
    ('Phase2-6\n规则挖掘', DIVIDER_GREEN),
    ('rules.json\n+ 自然语言规则', BG_WHITE),
]
fx_positions = [0.7, 2.1, 3.7, 5.1, 6.8, 8.2, 10.0]
for i, ((text, color), xp) in enumerate(zip(flow_items, fx_positions)):
    add_rect(slide, xp, 4.8, 1.4, 1.1, color)
    tc = TEXT_WHITE if color in [DIVIDER_BLUE, DIVIDER_ORANGE, DIVIDER_GREEN] else TEXT_DARK
    add_textbox(slide, xp + 0.05, 4.85, 1.3, 1.0, text, font_size=10, color=tc, bold=True, alignment=PP_ALIGN.CENTER)
    if i < len(flow_items) - 1:
        add_arrow_right(slide, xp + 1.45, 5.25, 0.45, 0.2, TEXT_GRAY)

add_textbox(slide, 0.7, 6.1, 11.9, 0.7,
            '关键价值：标注工具只做"好不好"的判断；规则挖掘工具自动从"不好的"里面挖规则；桥接让两者无缝衔接。',
            font_size=12, color=TEXT_DARK)


# ==================== 第 6 页：章节分隔 02 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide, '02', '完整流水线设计', 'Phase1-6 端到端架构')


# ==================== 第 7 页：全局架构图 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '黄金数据标注 × 冷数据规则提取 联调架构')

phases = [
    ('Phase1\n黄金数据标注', DIVIDER_BLUE,
     '无标注轨迹 → 全局理解 → 直觉评判\n输出: golden_output.jsonl'),
    ('桥接筛选', DIVIDER_ORANGE,
     '过滤 result∈{失败, 部分通过}\n输出: badcase_traces/'),
    ('Phase2\n链路归纳', PHASE_BLUE,
     '归纳正确链路 + 缺失检查点\n输出: phase2output.json'),
    ('Phase3\n类别归纳', PHASE_GREEN,
     '归纳 badcase 类别 + 二进制检查点\n输出: phase3output.json'),
    ('Phase4\n特征提取', PHASE_ORANGE,
     'LLM 逐条判定 0/1/NA\n输出: phase4output.json'),
    ('Phase5\n规则挖掘', PHASE_PURPLE,
     'IF-THEN + Skill 归因 + Few-Shot\n输出: rules.json'),
    ('Phase6\n规则转换', DIVIDER_GREEN,
     '自然语言化 + 三维排序筛选\n输出: rules_natural_language.txt\nrules_ranked.json'),
]

y = 1.4
for i, (title, color, desc) in enumerate(phases):
    add_rect(slide, 0.4, y, 12.5, 0.75, BG_LIGHT)
    tag = add_rect(slide, 0.4, y, 1.8, 0.75, color)
    add_textbox(slide, 0.5, y + 0.05, 1.6, 0.65, title,
                font_size=11, color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 2.3, y + 0.1, 10.4, 0.55, desc,
                font_size=11, color=TEXT_DARK)
    if i < len(phases) - 1:
        add_arrow_down(slide, 6.4, y + 0.75, 0.3, 0.12, TEXT_GRAY)
    y += 0.87

# 关键信息条
add_rect(slide, 0.4, 6.4, 12.5, 0.8, ACCENT_GREEN, DIVIDER_GREEN)
add_textbox(slide, 0.6, 6.45, 12, 0.7,
            '输入假设：Phase1 仅需"通过/失败"直觉标注  |  桥接严格过滤通过的轨迹  |  冷数据输入全部是 badcase\n'
            '最终产出：rules.json（可执行）+ rules_natural_language.txt（给 checker）+ rules_ranked.json（排序后）',
            font_size=11, color=TEXT_DARK, alignment=PP_ALIGN.CENTER)


# ==================== 第 8 页：Phase1 黄金数据标注 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, 'Phase1：直觉式黄金数据标注')

# 左栏
add_rect(slide, 0.5, 1.5, 6.0, 5.0, ACCENT_BLUE, DIVIDER_BLUE)
add_textbox(slide, 0.7, 1.6, 5.6, 0.4, '两阶段流程', font_size=18, color=DIVIDER_BLUE, bold=True)
steps = [
    ('阶段一：全局场景理解',
     '读取全部无标注轨迹，LLM 总结业务场景、用户目标、常见陷阱\n输出: global_understanding.txt'),
    ('阶段二：逐条直觉评判',
     '每条轨迹给出 expected_behavior + result + reason\nresult ∈ {通过, 失败, 部分通过}'),
]
y = 2.2
for title, desc in steps:
    add_textbox(slide, 0.7, y, 5.6, 0.35, title, font_size=14, color=DIVIDER_BLUE, bold=True)
    add_textbox(slide, 0.7, y + 0.4, 5.6, 1.1, desc, font_size=11, color=TEXT_DARK)
    y += 1.8

# 右栏
add_rect(slide, 6.8, 1.5, 5.8, 5.0, BG_LIGHT)
add_textbox(slide, 7.0, 1.6, 5.4, 0.4, '关键设计', font_size=18, color=DIVIDER_BLUE, bold=True)
designs = [
    ('低成本标注', '人工仅需判断"合格/不合格"\n无需分析具体错误类别'),
    ('输出增强（核心改动）', '保留原始 conversation_id / script_id / history / max_turns / total_turns\n供下游 Phase2-4 读取完整对话'),
    ('可断点续跑', '支持全局理解缓存复用\n避免重复调用 LLM'),
]
y = 2.2
for title, desc in designs:
    add_textbox(slide, 7.0, y, 5.4, 0.35, title, font_size=14, color=TEXT_DARK, bold=True)
    add_textbox(slide, 7.0, y + 0.4, 5.4, 0.9, desc, font_size=11, color=TEXT_DARK)
    y += 1.5

add_bottom_bar(slide, '核心改动：generate_golden() 返回 dict 追加 5 个原始轨迹字段')


# ==================== 第 9 页：Phase2 链路归纳 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, 'Phase2：从 Badcase 反推正确链路')

# 左栏
add_rect(slide, 0.4, 1.5, 5.5, 5.0, ACCENT_GREEN, PHASE_BLUE)
add_textbox(slide, 0.6, 1.55, 5.1, 0.4, '批次迭代归纳', font_size=16, color=PHASE_BLUE, bold=True)
iter_items = [
    ('Batch 1 (10条)', 'LLM 从零归纳', '标准业务链路\n+ 初始检查点'),
    ('Batch 2+', 'LLM 审阅优化', '合并/新增检查点\n更新链路'),
]
y = 2.1
for i, (title, action, output) in enumerate(iter_items):
    add_rect(slide, 0.6, y, 5.1, 1.3, BG_WHITE, PHASE_BLUE)
    add_textbox(slide, 0.7, y + 0.05, 2.2, 0.3, title, font_size=12, color=PHASE_BLUE, bold=True)
    add_textbox(slide, 0.7, y + 0.4, 2.2, 0.5, action, font_size=10, color=TEXT_DARK)
    add_textbox(slide, 3.0, y + 0.2, 2.5, 0.9, '→ ' + output, font_size=10, color=TEXT_DARK)
    if i < 1:
        add_arrow_down(slide, 2.8, y + 1.3, 0.35, 0.15, TEXT_GRAY)
    y += 1.6

# 右栏
add_rect(slide, 6.3, 1.5, 6.3, 5.0, BG_LIGHT)
add_textbox(slide, 6.5, 1.55, 5.9, 0.4, '关键设计', font_size=16, color=PHASE_BLUE, bold=True)
designs = [
    ('强制约束而非注意事项',
     '❌ "注意金额精度"\n✅ "金额计算保留原始精度至两位小数，禁止截断/取整"\n'
     'Prompt 要求输出：触发条件 + 判定标准 + 违反动作'),
    ('默认不去重',
     '重复轨迹代表场景高频出现\n需重点关注，去重改为显式 --dedup'),
    ('格式适配器',
     '自动修复 LLM 输出格式漂移\n字符串 → 字典，缺失字段补全'),
]
y = 2.1
for title, desc in designs:
    add_textbox(slide, 6.5, y, 5.9, 0.35, title, font_size=13, color=TEXT_RED, bold=True)
    add_textbox(slide, 6.5, y + 0.4, 5.9, 1.1, desc, font_size=10, color=TEXT_DARK)
    y += 1.5

add_bottom_bar(slide, '实战效果：19 条 badcase → 8 个步骤 → 9 个缺失检查点')


# ==================== 第 10 页：Phase3 类别归纳 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, 'Phase3：从检查点归纳可判定类别')

# 左栏
add_rect(slide, 0.4, 1.5, 6.0, 5.0, ACCENT_GREEN, PHASE_GREEN)
add_textbox(slide, 0.6, 1.55, 5.6, 0.4, '两阶段归纳', font_size=16, color=PHASE_GREEN, bold=True)
steps = [
    ('Step 1: LLM 自由聚类',
     '不预设分类框架\n让类别名称从错误模式中自然生长'),
    ('Step 2: 审查校验',
     '① 可判定性：能否给出 True/False\n'
     '② 客观性：是否不含主观评价词\n'
     '③ 证据位置明确：审查人能否在轨迹中找到证据'),
]
y = 2.1
for title, desc in steps:
    add_textbox(slide, 0.6, y, 5.6, 0.35, title, font_size=14, color=PHASE_GREEN, bold=True)
    add_textbox(slide, 0.6, y + 0.4, 5.6, 1.4, desc, font_size=11, color=TEXT_DARK)
    y += 2.0

# 右栏
add_rect(slide, 6.8, 1.5, 5.8, 5.0, BG_LIGHT)
add_textbox(slide, 7.0, 1.55, 5.4, 0.4, '类别数量控制', font_size=16, color=PHASE_GREEN, bold=True)
add_textbox(slide, 7.0, 2.1, 5.4, 1.2,
            '新增参数：\n--min-categories / --max-categories\n'
            '按轨迹数量动态约束类别数',
            font_size=12, color=TEXT_DARK)
add_textbox(slide, 7.0, 3.4, 5.4, 0.35, '核心原则', font_size=14, color=PHASE_GREEN, bold=True)
add_textbox(slide, 7.0, 3.8, 5.4, 1.2,
            '类别数 ≈ 检查点数 / 2~3\n每类 2-3 个检查点\n每条规则至少 2 条支撑轨迹',
            font_size=12, color=TEXT_DARK)

add_bottom_bar(slide, '实战效果：9 个检查点 → 5 个类别 / 8 个二进制检查点 / 100% 轨迹覆盖')


# ==================== 第 11 页：Phase4 特征提取 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, 'Phase4：LLM 理解式特征提取')

# 上部矩阵
add_rect(slide, 0.4, 1.5, 8.5, 2.8, BG_LIGHT)
add_textbox(slide, 0.6, 1.55, 8, 0.35, '特征矩阵 (轨迹 × 检查点)', font_size=16, color=PHASE_ORANGE, bold=True)
matrix_text = (
    '              CHK001  CHK002  CHK003  CHK004  CHK005  CHK006  CHK007  CHK008\n'
    'auto-2618e5  │  0   │  0   │  1   │  0   │  0   │  0   │  NA  │  1   │\n'
    'auto-39f227  │  NA  │  1   │  0   │  0   │  0   │  1   │  NA  │  NA  │\n'
    'auto-5faa84  │  0   │  NA  │  1   │  1   │  NA  │  0   │  NA  │  1   │\n'
    'auto-df5ee3  │  NA  │  NA  │  NA  │  NA  │  NA  │  NA  │  0   │  1   │\n\n'
    '0 = 违反(False)    1 = 通过(True)    NA = 不适用'
)
add_textbox(slide, 0.7, 2.0, 7.8, 2.2, matrix_text, font_size=10, color=TEXT_DARK, font_name='Consolas')

# 右上判定过程
add_rect(slide, 9.2, 1.5, 3.7, 2.8, ACCENT_PINK, PHASE_ORANGE)
add_textbox(slide, 9.4, 1.55, 3.3, 0.35, 'LLM 判定过程', font_size=14, color=PHASE_ORANGE, bold=True)
add_textbox(slide, 9.4, 2.0, 3.3, 2.1,
            '输入:\n  ① 检查点定义\n     (description +\n      judgment_criteria)\n  ② 轨迹对话片段\n\n输出:\n  final: 0/1/NA\n  reason: 判定依据\n  confidence: 0~1',
            font_size=11, color=TEXT_DARK)

# 下部关键设计
add_rect(slide, 0.4, 4.6, 12.5, 2.5, BG_LIGHT)
add_textbox(slide, 0.6, 4.65, 12, 0.35, '关键设计', font_size=16, color=PHASE_ORANGE, bold=True)
p4_designs = [
    ('三值判定(0/1/NA)', 'NA 表示检查点与此轨迹无关，后续统计排除，不污染 fail 率'),
    ('批量判定模式', '一次 LLM 调用同时判定多个检查点；--no-batch 切换逐个判定'),
    ('缓存机制', '相同 checkpoint+轨迹不重复调用；定义变更自动失效'),
    ('轨迹格式配置化', 'trajectory_config.json：字段名/角色映射/证据规则自定义'),
]
y = 5.15
for i, (title, desc) in enumerate(p4_designs):
    x = 0.6 + (i % 2) * 6.2
    yy = y + (i // 2) * 1.0
    add_textbox(slide, x, yy, 5.8, 0.25, '▸ ' + title, font_size=12, color=PHASE_ORANGE, bold=True)
    add_textbox(slide, x, yy + 0.28, 5.8, 0.6, desc, font_size=10, color=TEXT_DARK)


# ==================== 第 12 页：Phase5 规则挖掘 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, 'Phase5：规则挖掘 + Skill 归因')

# 左栏流程
add_rect(slide, 0.4, 1.5, 5.8, 5.0, ACCENT_PURPLE, PHASE_PURPLE)
add_textbox(slide, 0.6, 1.55, 5.4, 0.4, '规则挖掘流程', font_size=16, color=PHASE_PURPLE, bold=True)
flow = [
    ('按类别聚合特征矩阵', '统计每个 checkpoint 的 fail 率'),
    ('筛选 top-K 检查点 + 典型轨迹', 'fail 率最高的 K 个 + 违反分数最高的轨迹'),
    ('LLM 规则归纳', 'error_reason + if_conditions + skill_attribution + few_shots'),
    ('后处理', 'skill 白名单校验 + if 条件补描述\n重试 3 次 + 降级兜底'),
]
y = 2.1
for i, (title, desc) in enumerate(flow):
    add_rect(slide, 0.6, y, 5.4, 1.0, BG_WHITE, PHASE_PURPLE)
    add_textbox(slide, 0.7, y + 0.03, 5.2, 0.3, f'{i+1}. {title}', font_size=12, color=PHASE_PURPLE, bold=True)
    add_textbox(slide, 0.7, y + 0.35, 5.2, 0.6, desc, font_size=10, color=TEXT_DARK)
    if i < 3:
        add_arrow_down(slide, 3.0, y + 1.0, 0.3, 0.12, TEXT_GRAY)
    y += 1.2

# 右栏 JSON 示例
add_rect(slide, 6.6, 1.5, 6.1, 5.0, BG_LIGHT)
add_textbox(slide, 6.8, 1.55, 5.7, 0.4, '产出规则示例 (JSON)', font_size=14, color=PHASE_PURPLE, bold=True)
rule_json = (
    '{\n'
    '  "id": "R004",\n'
    '  "error_category": "CAT004-资金调拨与操作授权缺失",\n'
    '  "error_reason": "agent 资金不足时未请求授权...",\n'
    '  "if_conditions": [\n'
    '    {"feature":"CHK006_final","op":"==","value":0}\n'
    '  ],\n'
    '  "skill_attribution": {\n'
    '    "top3": [\n'
    '      {"skill_name":"fund_planning_skill",\n'
    '       "problematic_rule":"缺少资金调拨前确认",\n'
    '       "confidence":0.9}\n'
    '    ]\n'
    '  },\n'
    '  "few_shots": {...},\n'
    '  "confidence": 0.90\n'
    '}'
)
add_textbox(slide, 6.8, 2.1, 5.7, 4.2, rule_json, font_size=10, color=DIVIDER_BLUE, font_name='Consolas')


# ==================== 第 13 页：Phase6 规则转换与排序 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, 'Phase6：规则语言化 + 三维排序筛选')

# 左栏
add_rect(slide, 0.5, 1.5, 6.0, 5.0, ACCENT_GREEN, DIVIDER_GREEN)
add_textbox(slide, 0.7, 1.6, 5.6, 0.4, '规则语言化', font_size=18, color=DIVIDER_GREEN, bold=True)
nl_lines = [
    ('目标：将 rules.json 转为自然语言规则文本', 12, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('', 6, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('格式："当…时，必须…，否则判定为…"', 12, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('', 6, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('用途：作为评估器（checker）的提示词输入', 12, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('', 6, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('示例：', 12, TEXT_DARK, True, PP_ALIGN.LEFT),
    ('当购买金额超出当前账户余额、需要跨账户调拨时，', 11, TEXT_GRAY, False, PP_ALIGN.LEFT),
    ('agent 必须在告知用户缺口并请求指示后，才执行转账；', 11, TEXT_GRAY, False, PP_ALIGN.LEFT),
    ('否则判定为"资金调拨与操作授权缺失"。', 11, TEXT_GRAY, False, PP_ALIGN.LEFT),
]
add_rich_textbox(slide, 0.7, 2.2, 5.6, 4.0, nl_lines)

# 右栏
add_rect(slide, 6.8, 1.5, 5.8, 5.0, ACCENT_BLUE, DIVIDER_BLUE)
add_textbox(slide, 7.0, 1.6, 5.4, 0.4, '三维排序筛选', font_size=18, color=DIVIDER_BLUE, bold=True)
rank_items = [
    ('合理性', 'LLM 对抗性审查\n规则逻辑自洽、IF 可判定、THEN 合理'),
    ('重要性', 'severity（高/中/低）+ 业务影响面'),
    ('出现频率', 'supporting_trajectories 数量 / 总轨迹数\n+ 该类别下 checkpoint fail 率'),
]
y = 2.2
for title, desc in rank_items:
    add_textbox(slide, 7.0, y, 5.4, 0.35, title, font_size=14, color=DIVIDER_BLUE, bold=True)
    add_textbox(slide, 7.0, y + 0.4, 5.4, 0.9, desc, font_size=11, color=TEXT_DARK)
    y += 1.4
add_textbox(slide, 7.0, 6.0, 5.4, 0.4,
            'score = 0.4×合理性 + 0.3×重要性 + 0.3×频率',
            font_size=12, color=TEXT_DARK, bold=True)


# ==================== 第 14 页：章节分隔 03 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide, '03', '方法论沉淀', '提示词工程与经验元体系')


# ==================== 第 15 页：提示词优化 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '提示词优化两大方法论')

# 左栏
add_rect(slide, 0.5, 1.5, 6.0, 5.0, ACCENT_GREEN, DIVIDER_GREEN)
add_textbox(slide, 0.7, 1.6, 5.6, 0.4, '系统理解 6 点（Phase2）', font_size=18, color=DIVIDER_GREEN, bold=True)
sys_points = [
    '① 系统概况：agent 的角色和能力边界',
    '② 常见场景：轨迹中典型用户场景',
    '③ 用户目标：用户通常想达成什么',
    '④ 常见转折：关键节点与转折',
    '⑤ 常见陷阱：agent 最易出错环节',
    '⑥ 系统缺陷模式：技术限制 vs 逻辑缺陷',
]
y = 2.2
for p in sys_points:
    add_textbox(slide, 0.7, y, 5.6, 0.35, p, font_size=12, color=TEXT_DARK)
    y += 0.55
add_textbox(slide, 0.7, 5.6, 5.6, 0.6,
            '内化思考指引，不要求 LLM 逐条输出',
            font_size=11, color=TEXT_GRAY)

# 右栏
add_rect(slide, 6.8, 1.5, 5.8, 5.0, ACCENT_BLUE, DIVIDER_BLUE)
add_textbox(slide, 7.0, 1.6, 5.4, 0.4, '聚焦最终状态原则（Phase3）', font_size=18, color=DIVIDER_BLUE, bold=True)
final_lines = [
    ('描述「最终应达成的正确状态」', 12, TEXT_DARK, True, PP_ALIGN.LEFT),
    ('而非「中间应执行什么动作」', 12, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('', 6, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('从原则层面概括，不陷入细节', 12, TEXT_DARK, True, PP_ALIGN.LEFT),
    ('', 6, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('聚焦安全、合规、用户体验', 12, TEXT_DARK, True, PP_ALIGN.LEFT),
    ('', 6, TEXT_DARK, False, PP_ALIGN.LEFT),
    ('用"确保…""在…前提下完成…"表述', 12, TEXT_DARK, True, PP_ALIGN.LEFT),
]
add_rich_textbox(slide, 7.0, 2.2, 5.4, 2.8, final_lines)

# 底部案例对比
add_rect(slide, 0.5, 6.5, 12.3, 0.7, BG_LIGHT)
add_textbox(slide, 0.7, 6.55, 11.9, 0.6,
            '示例对比：❌ "agent 是否清洗金额为纯数字"  →  ✅ "确保用户在清晰获知标准化金额的前提下予以确认"',
            font_size=12, color=TEXT_DARK, alignment=PP_ALIGN.CENTER)


# ==================== 第 16 页：经验元体系 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '经验元：评估体系的持续进化单元')

# 定义
add_rect(slide, 0.5, 1.4, 12.3, 1.0, ACCENT_BLUE, DIVIDER_BLUE)
add_textbox(slide, 0.7, 1.45, 11.9, 0.9,
            '经验元 = 智能体在执行任务过程中产生的"最小可复用经验单元"\n'
            '规则元、反思元、案例元三位一体，支撑评估体系持续进化',
            font_size=15, color=TEXT_DARK, alignment=PP_ALIGN.CENTER)

# 三类卡片
exp_cards = [
    ('规则元', PHASE_BLUE, ACCENT_BLUE,
     ['由冷启动规则挖掘工具自动产出',
      '格式：IF-THEN + Few-Shot + Skill 归因',
      '作用：评估器直接执行判定',
      '🔄 新数据积累后刷新规则']),
    ('反思元', PHASE_PURPLE, ACCENT_PURPLE,
     ['由评估器在判定过程中反思产生',
      '格式：错误模式 + 根因 + 改进建议',
      '作用：补充规则元未覆盖的边界',
      '🔄 周期性归纳沉淀为新规则']),
    ('案例元', PHASE_ORANGE, ACCENT_PINK,
     ['由 Few-Shot 总结为黄金语料',
      '格式：标准化正例/负例轨迹',
      '作用：作为评估器判定参考标准',
      '🔄 评估反馈更新黄金语料库']),
]
for i, (title, tag_color, bg_color, bullets) in enumerate(exp_cards):
    x = 0.5 + i * 4.2
    add_rect(slide, x, 2.7, 3.9, 4.0, bg_color, tag_color)
    tag = add_rect(slide, x, 2.7, 1.6, 0.45, tag_color)
    add_textbox(slide, x + 0.05, 2.72, 1.5, 0.4, title, font_size=16,
                color=TEXT_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    y = 3.3
    for bullet in bullets:
        add_textbox(slide, x + 0.15, y, 3.6, 0.55, '• ' + bullet, font_size=11, color=TEXT_DARK)
        y += 0.6


# ==================== 第 17 页：章节分隔 04 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide, '04', '实战案例与效果', '银行理财购买 Agent 端到端运行')


# ==================== 第 18 页：实战数据与规则产出 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '案例：银行理财购买 Agent')

# 上方输入产出
add_rect(slide, 0.4, 1.5, 5.8, 2.0, ACCENT_GREEN, DIVIDER_GREEN)
add_textbox(slide, 0.6, 1.55, 5.4, 0.4, '场景与输入', font_size=16, color=DIVIDER_GREEN, bold=True)
add_textbox(slide, 0.6, 2.05, 5.4, 1.3,
            '• 业务：工行理财购买对话 Agent\n'
            '• 组成：AgentRule + product_select_skill + fund_planning_skill\n'
            '• 输入：20 条无标注轨迹（0611v1）\n'
            '• Phase1 标注：1 通过 / 11 失败 / 8 部分通过',
            font_size=12, color=TEXT_DARK)

add_rect(slide, 6.6, 1.5, 6.1, 2.0, ACCENT_BLUE, DIVIDER_BLUE)
add_textbox(slide, 6.8, 1.55, 5.7, 0.4, '产出概览', font_size=16, color=DIVIDER_BLUE, bold=True)
add_textbox(slide, 6.8, 2.05, 5.7, 1.3,
            '• Phase2：8 个步骤 / 9 个缺失检查点\n'
            '• Phase3：5 个类别 / 8 个二进制检查点\n'
            '• Phase4：19 × 8 = 152 个判定\n'
            '• Phase5-6：5 条可部署规则（全部保留）',
            font_size=12, color=TEXT_DARK)

# 下方规则表
add_rect(slide, 0.4, 3.8, 12.3, 3.0, BG_LIGHT)
add_textbox(slide, 0.6, 3.9, 11.9, 0.4, '5 条规则排序结果', font_size=16, color=TEXT_DARK, bold=True)

rules_data = [
    ('#1', 'R003', '交易最终确认与修改确认缺失', 'CHK004=0 AND CHK005=0', 'AgentRule', '0.833'),
    ('#2', 'R004', '资金调拨与操作授权缺失', 'CHK006=0', 'fund_planning_skill', '0.833'),
    ('#3', 'R001', '身份与账户确认缺失', 'CHK001=0', 'AgentRule', '0.831'),
    ('#4', 'R002', '输入格式与名称匹配确认缺失', 'CHK002=0 AND CHK003=0', 'product_select_skill', '0.751'),
    ('#5', 'R005', '信息查询与系统异常处理不当', 'CHK008=0', 'AgentRule_added_v5', '0.695'),
]

# 表头
headers2 = ['排名', '规则', '错误类别', '触发条件', '首位 Skill', '综合得分']
xs2 = [0.6, 1.3, 2.4, 5.4, 8.6, 11.2]
for x, h in zip(xs2, headers2):
    add_textbox(slide, x, 4.35, 1.8, 0.3, h, font_size=10, color=DIVIDER_BLUE, bold=True)

y = 4.7
for rank, rid, cat, cond, skill, score in rules_data:
    add_textbox(slide, 0.6, y, 0.7, 0.3, rank, font_size=10, color=TEXT_DARK)
    add_textbox(slide, 1.3, y, 1.0, 0.3, rid, font_size=10, color=TEXT_DARK)
    add_textbox(slide, 2.4, y, 2.8, 0.3, cat, font_size=10, color=PHASE_PURPLE)
    add_textbox(slide, 5.4, y, 3.0, 0.3, cond, font_size=9, color=PHASE_BLUE, font_name='Consolas')
    add_textbox(slide, 8.6, y, 2.4, 0.3, skill, font_size=10, color=PHASE_ORANGE, bold=True)
    add_textbox(slide, 11.2, y, 1.0, 0.3, score, font_size=10, color=TEXT_DARK, alignment=PP_ALIGN.CENTER)
    y += 0.36

add_bottom_bar(slide, '🌟 关键发现：AgentRule（全局规则）在 5 条规则中 4 条进入 Top3，底层约束缺失是系统性根因')


# ==================== 第 19 页：章节分隔 05 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_divider(slide, '05', '核心价值与展望', '从工具到体系的进化路径')


# ==================== 第 20 页：核心价值提炼 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '核心价值提炼')

# 左栏
add_rect(slide, 0.5, 1.5, 6.0, 5.0, ACCENT_GREEN, DIVIDER_GREEN)
add_textbox(slide, 0.7, 1.6, 5.6, 0.4, '对业务的价值', font_size=18, color=DIVIDER_GREEN, bold=True)
biz_values = [
    ('缩短规则产出周期', '从 2-4 周 → 2-4 小时'),
    ('降低标注成本', '仅需直觉式通过/失败标注'),
    ('提升规则可执行性', '直接输出 IF-THEN + Skill 归因'),
    ('定位优化方向', '规则直接关联 Skill 文件，优化有据可依'),
]
y = 2.2
for title, desc in biz_values:
    add_textbox(slide, 0.7, y, 5.6, 0.35, title, font_size=14, color=DIVIDER_GREEN, bold=True)
    add_textbox(slide, 0.7, y + 0.4, 5.6, 0.5, desc, font_size=12, color=TEXT_DARK)
    y += 1.0

# 右栏
add_rect(slide, 6.8, 1.5, 5.8, 5.0, ACCENT_BLUE, DIVIDER_BLUE)
add_textbox(slide, 7.0, 1.6, 5.4, 0.4, '对技术的价值', font_size=18, color=DIVIDER_BLUE, bold=True)
tech_values = [
    ('端到端自动化', 'Phase1-6 一键运行'),
    ('工程健壮', '断点续跑 + 缓存 + 重试降级'),
    ('方法论可复现', '系统理解 6 点 + 最终状态原则'),
    ('经验元体系', '规则元·反思元·案例元持续进化'),
]
y = 2.2
for title, desc in tech_values:
    add_textbox(slide, 7.0, y, 5.4, 0.35, title, font_size=14, color=DIVIDER_BLUE, bold=True)
    add_textbox(slide, 7.0, y + 0.4, 5.4, 0.5, desc, font_size=12, color=TEXT_DARK)
    y += 1.0


# ==================== 第 21 页：未来方向 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_title_bar(slide, '未来方向')

futures = [
    ('🔮 海量数据高效归纳', ACCENT_BLUE, DIVIDER_BLUE,
     '1000+ 轨迹时引入 BERTopic 等主题建模工具\n替代纯 LLM 聚类，降低成本、提升速度'),
    ('📊 主动学习闭环', ACCENT_GREEN, DIVIDER_GREEN,
     '评估器判定结果反哺 Phase2-6\n积累新数据 → 规则刷新 → 评估更精准'),
    ('🌐 跨领域迁移', ACCENT_PINK, DIVIDER_PURPLE,
     '同一框架复用于客服、医疗、保险等\n对话 Agent 场景'),
]
for i, (title, bg, accent, desc) in enumerate(futures):
    x = 0.5 + i * 4.2
    add_rect(slide, x, 1.5, 3.9, 3.5, bg, accent)
    add_rect(slide, x, 1.5, 3.9, 0.08, accent, corner=False)
    add_textbox(slide, x + 0.2, 1.7, 3.5, 0.4, title, font_size=16, color=accent, bold=True)
    add_textbox(slide, x + 0.2, 2.2, 3.5, 2.5, desc, font_size=13, color=TEXT_DARK)

# 底部闭环
add_rect(slide, 0.5, 5.3, 12.3, 1.8, BG_LIGHT)
add_textbox(slide, 0.7, 5.4, 11.9, 0.4, '从冷启动到持续进化的闭环', font_size=16, color=DIVIDER_BLUE, bold=True)
add_textbox(slide, 0.7, 5.9, 11.9, 1.0,
            '初始 50-100 条标注 → 冷启动规则挖掘 → 评估器上线 → 反思积累经验元 → '
            '运行积累新数据（1000+）→ BERTopic 高效归纳 → 规则刷新 → 评估器升级 → 循环持续进化',
            font_size=13, color=TEXT_DARK, alignment=PP_ALIGN.CENTER)


# ==================== 第 22 页：结尾页 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_plain_rect(slide, 0, 6.9, 13.333, 0.6, DIVIDER_GREEN)

add_textbox(slide, 1.5, 2.4, 10, 1.2, '从 30 条 Badcase 到可部署评估规则',
            font_size=40, color=TEXT_DARK, bold=True, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 3.6, 10, 0.8,
            '黄金数据标注 × 冷数据规则提取 联调方案',
            font_size=24, color=DIVIDER_BLUE, alignment=PP_ALIGN.CENTER)
add_textbox(slide, 1.5, 4.6, 10, 0.8,
            '让 LLM Agent 的评估规则，像数据一样自动生长',
            font_size=20, color=DIVIDER_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

add_textbox(slide, 4, 6.0, 5, 0.4, '2026.06', font_size=14,
            color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)


# ==================== 保存 ====================
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'kimiv4ppt.pptx')
prs.save(output_path)
print(f'PPT已保存: {output_path}')
